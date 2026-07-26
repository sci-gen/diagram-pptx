import inspect
import json
from pathlib import Path
from zipfile import ZipFile

import pytest
from pptx import Presentation
from pptx.util import Inches

from diagram_pptx import (
    PartialModelMutationError,
    compile_diagram,
    parse_mermaid,
    render_mermaid,
    resolve_diagram_bounds,
    sample_colormap,
)
from diagram_pptx.scene import SceneConnector, SceneShape


def test_partial_model_mutation_is_rejected() -> None:
    document = parse_mermaid("flowchart LR\nA --> B\nclick A https://example.com")
    document.model.nodes["A"].label = "Changed"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(PartialModelMutationError):
        compile_diagram(
            document,
            slide=slide,
            bounds=(1, 1, 8, 4),
            backend="native",
        )


def test_grouped_output_survives_save_and_reopen(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    document = parse_mermaid("flowchart LR\nA[One] --> B[Two]")
    result = compile_diagram(
        document,
        slide=slide,
        bounds=(1, 1, 8, 4),
        backend="native",
    )
    result.group_shape.width = 7_000_000
    result.group_shape.height = 3_000_000
    output = tmp_path / "grouped.pptx"
    presentation.save(output)

    reopened = Presentation(output)
    group = reopened.slides[0].shapes[0]
    assert len(reopened.slides[0].shapes) == 1
    assert group.name == "diagram:flowchart"
    assert group.width == 7_000_000
    assert group.height == 3_000_000
    assert len(group.shapes) == len(result.shapes)
    assert {"A", "B"} <= set(result.element_shapes)

    with ZipFile(output) as archive:
        xml = archive.read("ppt/slides/slide1.xml").decode()
    assert "<p:grpSp>" in xml
    assert "<p:sp>" in xml
    assert "<p:cxnSp>" in xml
    assert "<a:blip" not in xml


def test_rgba_and_opacity_are_written_as_drawingml_alpha(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    document = parse_mermaid("flowchart LR\nA[One]")
    document.model.nodes["A"].style.fill = "#33669980"
    document.model.nodes["A"].style.opacity = 0.5
    compile_diagram(
        document,
        slide=slide,
        bounds=(1, 1, 4, 2),
        backend="native",
    )
    output = tmp_path / "alpha.pptx"
    presentation.save(output)

    with ZipFile(output) as archive:
        xml = archive.read("ppt/slides/slide1.xml").decode()
    assert '<a:alpha val="25098"/>' in xml


def test_sequence_messages_remain_on_event_rows() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        parse_mermaid("sequenceDiagram\nparticipant A\nparticipant B\nA->>B: Hello"),
        slide=slide,
        bounds=(1, 1, 8, 4),
        backend="native",
    )

    messages = [shape for shape in result.connectors if "sequence.message" in shape.name]
    assert len(messages) == 1
    assert messages[0].begin_y == messages[0].end_y


def test_native_is_default_and_style_preset_is_independent() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        parse_mermaid("flowchart LR\nA[One] --> B{Two}"),
        slide=slide,
        bounds=(1, 1, 8, 4),
        style="official",
        group=False,
    )
    node = next(
        item
        for item in result.scene.elements
        if isinstance(item, SceneShape) and item.semantic_id == "A"
    )

    assert result.backend_used == "native"
    assert node.style.fill == "#ECECFF"
    assert node.style.line == "#9370DB"


def test_numeric_colormap_positions_override_node_and_edge_channels() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        parse_mermaid("flowchart LR\nA[One] --> B{Two}"),
        slide=slide,
        bounds=(1, 1, 8, 4),
        colors={
            "name": "jet",
            "fill": 0.65,
            "line": 0.9,
            "edge": 0.75,
            "text": 0.2,
            "decision_fill": 0.45,
        },
        group=False,
    )
    nodes = {
        item.semantic_id: item for item in result.scene.elements if isinstance(item, SceneShape)
    }
    connector = next(item for item in result.scene.elements if isinstance(item, SceneConnector))

    assert nodes["A"].style.fill == sample_colormap("jet", 0.65)
    assert nodes["B"].style.fill == sample_colormap("jet", 0.45)
    assert nodes["A"].style.line == sample_colormap("jet", 0.9)
    assert nodes["A"].style.text == sample_colormap("jet", 0.2)
    assert connector.style.line == sample_colormap("jet", 0.75)


def test_global_label_background_resolves_theme_tokens_and_element_override_wins() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        parse_mermaid("flowchart LR\nA[One] -->|valid| B[Two]"),
        slide=slide,
        bounds=(1, 1, 8, 4),
        theme={"palette": {"canvas": "#243447"}},
        label_background="canvas",
        group=False,
    )
    connector = next(item for item in result.scene.elements if isinstance(item, SceneConnector))

    assert connector.style.label_fill == "#243447"

    overridden = compile_diagram(
        parse_mermaid("flowchart LR\nA[One] -->|valid| B[Two]"),
        slide=presentation.slides.add_slide(presentation.slide_layouts[6]),
        bounds=(1, 1, 8, 4),
        label_background="#243447",
        style_overrides={"edge-0": {"label_fill": "#F3F0E8"}},
        group=False,
    )
    overridden_connector = next(
        item for item in overridden.scene.elements if isinstance(item, SceneConnector)
    )

    assert overridden_connector.style.label_fill == "#F3F0E8"


def test_primary_secondary_colors_choose_text_contrast_per_node() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        parse_mermaid("flowchart LR\nA[Dark] --> B{Light}"),
        slide=slide,
        bounds=(1, 1, 8, 4),
        colors={"name": "magma", "primary": 0.05, "secondary": 0.95},
        group=False,
    )
    nodes = {
        item.semantic_id: item for item in result.scene.elements if isinstance(item, SceneShape)
    }
    connector = next(item for item in result.scene.elements if isinstance(item, SceneConnector))

    assert nodes["A"].style.fill == sample_colormap("magma", 0.05)
    assert nodes["A"].style.text == "#FFFFFF"
    assert nodes["B"].style.fill == sample_colormap("magma", 0.95)
    assert nodes["B"].style.text == "#111827"
    assert connector.style.line == sample_colormap("magma", 0.95)


def test_json_like_arguments_and_result_support_external_tool_adapters() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    arguments = {
        "source": "flowchart LR\nA[One] --> B[Two]",
        "position": "left",
        "style": "official",
        "colors": {"name": "viridis", "primary": 0.8, "secondary": 0.2},
    }
    result = render_mermaid(slide=slide, **arguments)

    assert result.backend_used == "native"
    assert result.group_shape is not None
    summary = result.to_dict()
    assert summary["diagram_kind"] == "flowchart"
    assert summary["grouped"] is True
    assert summary["element_ids"]
    assert json.loads(json.dumps(summary)) == summary


def test_public_signature_uses_one_colors_object_and_keeps_alpha_compatibility() -> None:
    parameters = inspect.signature(render_mermaid).parameters

    assert "colors" in parameters
    assert "colormap" not in parameters
    assert "primary" not in parameters
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        "flowchart LR\nA --> B",
        slide=slide,
        colormap="jet",
        primary=0.8,
        secondary=0.2,
    )
    assert result.group_shape is not None


def test_colors_object_cannot_be_mixed_with_legacy_color_arguments() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(TypeError, match="Use colors"):
        render_mermaid(
            "flowchart LR\nA --> B",
            slide=slide,
            colors={"name": "jet"},
            colormap="viridis",
        )


def test_runtime_timeout_must_be_positive() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(ValueError, match="timeout must be greater"):
        render_mermaid("flowchart LR\nA --> B", slide=slide, timeout=0)


def test_position_presets_and_relative_bounds_use_slide_coordinates() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)
    left_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    right_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    relative_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    document = parse_mermaid("flowchart LR\nA[One] --> B[Two]")

    left = compile_diagram(document, slide=left_slide, position="left")
    right = compile_diagram(document, slide=right_slide, position="right")
    relative = compile_diagram(
        document,
        slide=relative_slide,
        relative_bounds=(0.1, 0.2, 0.3, 0.4),
    )

    slide_center = presentation.slide_width / 2
    assert left.group_shape.left + left.group_shape.width / 2 < slide_center
    assert right.group_shape.left + right.group_shape.width / 2 > slide_center
    assert relative.group_shape.left >= presentation.slide_width * 0.1
    assert relative.group_shape.top >= presentation.slide_height * 0.2
    assert relative.group_shape.left + relative.group_shape.width <= presentation.slide_width * 0.4
    assert relative.group_shape.top + relative.group_shape.height <= presentation.slide_height * 0.6


def test_multiple_diagrams_can_share_one_slide_without_overlapping_presets() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    document = parse_mermaid("flowchart LR\nA[One] --> B[Two]")

    left = compile_diagram(document, slide=slide, position="left")
    right = compile_diagram(document, slide=slide, position="right")

    assert len(slide.shapes) == 2
    assert left.group_shape.left + left.group_shape.width < right.group_shape.left


def test_diagram_bounds_are_mutually_exclusive_and_default_to_full() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    assert resolve_diagram_bounds(slide, position="left") == pytest.approx((0.5, 0.6, 4.2, 6.3))
    assert resolve_diagram_bounds(slide, position="right") == pytest.approx((5.3, 0.6, 4.2, 6.3))
    assert resolve_diagram_bounds(slide, position="top") == pytest.approx((0.5, 0.6, 9.0, 2.85))
    assert resolve_diagram_bounds(slide, position="bottom") == pytest.approx((0.5, 4.05, 9.0, 2.85))
    assert resolve_diagram_bounds(slide) == pytest.approx((0.5, 0.375, 9.0, 6.75))
    with pytest.raises(ValueError):
        resolve_diagram_bounds(
            slide,
            bounds=(1, 1, 4, 3),
            relative_bounds=(0, 0, 0.5, 0.5),
        )
    with pytest.raises(ValueError):
        resolve_diagram_bounds(slide, relative_bounds=(0.8, 0.8, 0.3, 0.3))
