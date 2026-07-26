from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from diagram_pptx import render_mermaid


def test_renders_only_native_shapes(tmp_path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        """
        flowchart LR
            A[Receive] --> B{Valid?}
            B -->|Yes| C[Process]
            B -.->|No| D[Reject]
        """,
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
    )
    output = tmp_path / "native.pptx"
    presentation.save(output)

    assert {"A", "B", "C", "D"} <= set(result.node_shapes)
    assert len(result.connectors) >= 3
    assert len(result.edge_label_shapes) == 2
    assert all(
        result.node_shapes[node_id].name.startswith("diagram:flowchart:node.")
        and result.node_shapes[node_id].name.endswith(f":{node_id}")
        for node_id in ("A", "B", "C", "D")
    )

    reopened = Presentation(output)
    assert len(reopened.slides[0].shapes) == 1
    assert reopened.slides[0].shapes[0].shape_type == 6
    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<p:grpSp>" in slide_xml
    assert "<p:cxnSp>" in slide_xml
    assert "tailEnd" in slide_xml
    assert "<a:blip" not in slide_xml


def test_lr_native_connectors_use_shape_specific_connection_sites() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        "flowchart LR\nA[Receive] --> B{Valid?}\n",
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )

    assert len(result.connectors) == 1
    connector_xml = result.connectors[0]._element.xml
    # Rectangle right side -> diamond left side.
    assert 'idx="3"' in connector_xml
    assert 'idx="1"' in connector_xml


def test_native_shapes_disable_theme_effects_and_use_medium_arrows() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        "flowchart LR\nA[Receive] -->|OK| B[Store]\n",
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )

    assert all(
        '<a:effectRef idx="0">' in shape._element.xml
        for shape in result.shapes
        if "<p:style>" in shape._element.xml
    )
    connector_xml = result.connectors[-1]._element.xml
    assert '<a:tailEnd type="triangle" w="med" len="med"/>' in connector_xml
    assert result.edge_label_shapes[0].width < Inches(1.0)


def test_sequence_message_label_is_above_the_line_without_an_opaque_fill() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        """
        sequenceDiagram
            participant A
            participant B
            A->>B: Send
        """,
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )

    line = next(
        connector for connector in result.connectors if "sequence.message" in connector.name
    )
    label = result.edge_label_shapes[0]
    assert label.top + label.height < line.top
    assert label.fill.type == MSO_FILL_TYPE.BACKGROUND


def test_short_connector_label_moves_above_the_line_without_hiding_the_arrow() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        "flowchart LR\nA[Start] -->|complete| B[Done]\n",
        slide=slide,
        bounds=(0.7, 0.7, 3.0, 2.0),
        backend="native",
        group=False,
    )

    line = result.connectors[-1]
    label = result.edge_label_shapes[0]
    assert label.top + label.height < line.top
    assert label.fill.type == MSO_FILL_TYPE.BACKGROUND
    assert '<a:tailEnd type="triangle" w="med" len="med"/>' in line._element.xml


def test_explicit_label_background_applies_to_detached_powerpoint_labels() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        "flowchart LR\nA[Start] -->|complete| B[Done]\n",
        slide=slide,
        bounds=(0.7, 0.7, 3.0, 2.0),
        label_background="#243447",
        group=False,
    )

    label = result.edge_label_shapes[0]
    assert label.fill.type == MSO_FILL_TYPE.SOLID
    assert str(label.fill.fore_color.rgb) == "243447"
    assert str(label.text_frame.paragraphs[0].runs[0].font.color.rgb) == "FFFFFF"


def test_bent_vertical_label_stays_on_line_while_short_horizontal_labels_detach() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        """
        stateDiagram-v2
            direction LR
            Idle --> Workflow : start
            state Workflow {
                [*] --> Validate
                Validate --> Execute : valid
                Execute --> [*]
            }
            Workflow --> Idle : cancel
            Workflow --> Done : complete
        """,
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )
    labels = {shape.text: shape for shape in result.edge_label_shapes}

    assert labels["valid"].fill.type == MSO_FILL_TYPE.SOLID
    assert labels["cancel"].fill.type == MSO_FILL_TYPE.BACKGROUND
    assert labels["complete"].fill.type == MSO_FILL_TYPE.BACKGROUND


def test_visual_container_uses_a_subtle_corner_radius() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        """
        flowchart LR
            subgraph G[Group]
                A --> B
            end
        """,
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )

    outline = result.group_shapes[0]
    assert outline.adjustments[0] == 0.06
    assert 'fmla="val 6000"' in outline._element.xml


def test_native_font_size_scales_with_fitted_diagram() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        "flowchart LR\nA[Receive] --> B[Store]\n",
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )

    run = result.node_shapes["A"].text_frame.paragraphs[0].runs[0]
    assert 15.0 < run.font.size.pt <= 28.0


def test_native_sequence_fragment_uses_square_frame() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    render_mermaid(
        """
        sequenceDiagram
            participant A
            participant B
            alt Ready
                A->>B: Send
            else Retry
                B-->>A: Wait
            end
        """,
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="native",
        group=False,
    )

    fragment = next(
        shape
        for shape in slide.shapes
        if "sequence.fragment.alt" in shape.name and not shape.name.endswith(":label")
    )
    assert fragment.auto_shape_type == MSO_SHAPE.RECTANGLE


def test_native_actor_and_composite_state_create_nested_groups() -> None:
    presentation = Presentation()
    actor_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    actor_result = render_mermaid(
        """
        sequenceDiagram
            actor User
            participant API
            User->>API: Request
        """,
        slide=actor_slide,
        group=True,
    )

    assert [shape.name for shape in actor_result.nested_group_shapes] == [
        "diagram:sequence:actor:User"
    ]
    assert actor_result.group_shape is not None

    state_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    state_result = render_mermaid(
        """
        stateDiagram-v2
            state Workflow {
                [*] --> Idle
                Idle --> Busy
            }
            Outside --> Workflow
            Workflow --> Outside
        """,
        slide=state_slide,
        group=True,
    )

    assert [shape.name for shape in state_result.nested_group_shapes] == [
        "diagram:state:composite:Workflow"
    ]
    assert state_result.group_shape is not None
    assert any(
        "idx=" in connector._element.xml
        for connector in state_result.connectors
        if "transition" in connector.name
    )


def test_nested_composite_states_create_hierarchical_groups() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = render_mermaid(
        """
        stateDiagram-v2
            state Outer {
                state Inner {
                    [*] --> Work
                    Work --> [*]
                }
            }
        """,
        slide=slide,
        group=True,
    )

    assert {shape.name for shape in result.nested_group_shapes} == {
        "diagram:state:composite:Outer",
        "diagram:state:composite:Inner",
    }
    assert result.group_shape is not None
