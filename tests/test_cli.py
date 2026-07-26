import json
from importlib.metadata import version
from pathlib import Path

from pptx import Presentation

import diagram_pptx
from diagram_pptx.cli import _package_version, main


def test_public_versions_match_distribution_metadata() -> None:
    expected = version("diagram-pptx")
    assert expected == "0.1.0b1"
    assert diagram_pptx.__version__ == expected
    assert _package_version() == expected


def test_cli_renders_a_presentation(tmp_path: Path) -> None:
    source = tmp_path / "diagram.mmd"
    output = tmp_path / "diagram.pptx"
    source.write_text("flowchart LR\nA[Start] -->|done| B[Finish]", encoding="utf-8")

    status = main(
        [
            "render",
            str(source),
            str(output),
            "--title",
            "Example",
            "--style",
            "official",
            "--colormap",
            "viridis",
            "--primary",
            "0.82",
            "--secondary",
            "0.20",
            "--label-background",
            "#F3F0E8",
        ]
    )

    assert status == 0
    assert output.exists()
    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 2
    assert presentation.slides[0].shapes[1].shape_type == 6


def test_cli_accepts_normalized_relative_bounds(tmp_path: Path) -> None:
    source = tmp_path / "diagram.mmd"
    output = tmp_path / "diagram.pptx"
    source.write_text("flowchart LR\nA[Start] --> B[Finish]", encoding="utf-8")

    status = main(
        [
            "render",
            str(source),
            str(output),
            "--relative-bounds",
            "0.05,0.10,0.40,0.80",
        ]
    )

    assert status == 0
    presentation = Presentation(output)
    group = presentation.slides[0].shapes[0]
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    assert group.left >= int(slide_width * 0.05)
    assert group.top >= int(slide_height * 0.10)
    assert group.left + group.width <= int(slide_width * 0.45) + 2
    assert group.top + group.height <= int(slide_height * 0.90) + 2


def test_cli_inspect_reports_model_coverage(
    tmp_path: Path,
    capsys,
) -> None:
    source = tmp_path / "partial.mmd"
    source.write_text(
        "flowchart LR\nA --> B\nclick A https://example.com",
        encoding="utf-8",
    )

    status = main(["inspect", str(source), "--json"])
    payload = json.loads(capsys.readouterr().out)

    assert status == 0
    assert payload["kind"] == "flowchart"
    assert payload["modeling_rate"] < 1
    assert payload["required_backend"] == "official"


def test_cli_doctor_reports_image_export(capsys) -> None:
    status = main(["doctor", "--json"])
    payload = json.loads(capsys.readouterr().out)

    assert status == 0
    assert payload["image_export"]["svg"] is True
    assert payload["image_export"]["png"] is True
    assert payload["image_export"]["jpeg"] is True


def test_cli_support_reports_all_registered_mermaid_families(capsys) -> None:
    status = main(["support", "--json"])
    payload = json.loads(capsys.readouterr().out)

    assert status == 0
    assert payload["mermaid_version"] == "11.16.0"
    assert payload["summary"] == {
        "registered": 31,
        "official": 31,
        "typed_model": 5,
        "native_backend": 5,
    }
    assert {row["kind"] for row in payload["families"]} >= {
        "flowchart",
        "sequence",
        "swimlanes",
        "eventmodeling",
        "railroad",
    }
