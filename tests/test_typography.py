from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

from diagram_pptx import (
    DiagramCompiler,
    DiagramSettings,
    DiagramTheme,
    ElementStyle,
    FontSize,
    TypographySettings,
    parse_mermaid,
)


def _font_points(shape) -> float:
    return shape.text_frame.paragraphs[0].runs[0].font.size.pt


def test_font_size_units_resolve_to_powerpoint_points() -> None:
    assert FontSize.pt(12).resolve(slide_height_points=360) == 12
    assert FontSize.px(16).resolve(slide_height_points=360) == 12
    assert FontSize.slide_height(0.04).resolve(slide_height_points=360) == 14.4
    assert FontSize.parse("2.5%sh").resolve(slide_height_points=400) == 10


def test_typography_settings_round_trip_with_default_point_unit() -> None:
    settings = DiagramSettings(
        typography=TypographySettings(
            font_family="Aptos",
            japanese_font_family="Yu Gothic",
            node=18,
            edge="16px",
            group=FontSize.slide_height(0.025),
        )
    )

    assert settings.to_dict()["typography"]["node"] == 18
    assert settings.to_dict()["typography"]["edge"] == "16px"
    assert settings.to_dict()["typography"]["group"] == "2.5%sh"
    assert settings.to_dict()["typography"]["font_family"] == "Aptos"
    assert settings.to_dict()["typography"]["japanese_font_family"] == "Yu Gothic"
    assert DiagramSettings.from_dict(settings.to_dict()).to_dict() == settings.to_dict()


def test_manual_font_precedence_and_slide_relative_element_override() -> None:
    presentation = Presentation()
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    document = parse_mermaid("flowchart LR\nA[Alpha] --> B[Beta]\n")
    document.model.nodes["A"].style.font_size = FontSize.slide_height(0.04)

    result = DiagramCompiler(
        DiagramSettings(typography=TypographySettings(node=16, fit="fit"))
    ).compile(
        document,
        slide=slide,
        bounds=(0.7, 0.7, 8.0, 3.5),
        theme=DiagramTheme(roles={"node.default": ElementStyle(font_size=18)}),
        style_overrides={"B": ElementStyle(font_size="20pt")},
        group=False,
    )

    # Element source style > role theme; compile override is final.
    assert _font_points(result.node_shapes["A"]) == 14.4
    assert _font_points(result.node_shapes["B"]) == 20
    assert result.node_shapes["A"].text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def test_fit_none_keeps_manual_size_without_powerpoint_autofit() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = DiagramCompiler(
        DiagramSettings(typography=TypographySettings(node="19pt", fit="none"))
    ).render_mermaid(
        "flowchart LR\nA[Alpha] --> B[Beta]\n",
        slide=slide,
        group=False,
    )

    assert _font_points(result.node_shapes["A"]) == 19
    assert result.node_shapes["A"].text_frame.auto_size == MSO_AUTO_SIZE.NONE


def test_absolute_points_do_not_change_with_diagram_bounds() -> None:
    presentation = Presentation()
    compiler = DiagramCompiler(DiagramSettings(typography=TypographySettings(node=17, fit="none")))
    sizes = []
    for bounds in ((0.5, 0.5, 3.0, 2.0), (0.5, 0.5, 10.0, 5.5)):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        result = compiler.render_mermaid(
            "flowchart LR\nA[Alpha] --> B[Beta]\n",
            slide=slide,
            bounds=bounds,
            group=False,
        )
        sizes.append(_font_points(result.node_shapes["A"]))

    assert sizes == [17, 17]


def test_japanese_text_defaults_to_yu_gothic_and_sets_east_asian_typeface() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = DiagramCompiler().render_mermaid(
        "flowchart LR\nA[申請を確認] --> B[承認]\n",
        slide=slide,
        group=False,
    )

    run = result.node_shapes["A"].text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Yu Gothic"
    assert '<a:ea typeface="Yu Gothic"/>' in run._r.xml


def test_japanese_font_and_mixed_element_sizes_are_instance_configurable() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    compiler = DiagramCompiler(
        DiagramSettings(
            typography=TypographySettings(
                japanese_font_family="BIZ UDPGothic",
                fit="fit",
            )
        )
    )
    result = compiler.render_mermaid(
        "flowchart LR\nA[開始] --> B[確認] --> C[完了]\n",
        slide=slide,
        bounds=(0.4, 0.8, 12.4, 5.5),
        style_overrides={
            "A": ElementStyle(font_size="14pt"),
            "B": ElementStyle(font_size="22pt"),
            "C": ElementStyle(font_size="30pt"),
        },
        group=False,
    )

    assert [_font_points(result.node_shapes[key]) for key in ("A", "B", "C")] == [
        14,
        22,
        30,
    ]
    assert {
        result.node_shapes[key].text_frame.paragraphs[0].runs[0].font.name
        for key in ("A", "B", "C")
    } == {"BIZ UDPGothic"}


def test_fit_pre_shrinks_long_text_before_powerpoint_opens() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = DiagramCompiler(
        DiagramSettings(typography=TypographySettings(node="28pt", fit="fit"))
    ).render_mermaid(
        "flowchart LR\nA[This label is deliberately much longer than its node]\n",
        slide=slide,
        bounds=(0.5, 0.5, 2.4, 1.2),
        group=False,
    )

    shape = result.node_shapes["A"]
    assert 9 <= _font_points(shape) < 28
    assert shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
