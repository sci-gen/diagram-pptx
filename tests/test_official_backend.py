import shutil
from pathlib import Path
from zipfile import ZipFile

import pytest
from pptx import Presentation
from test_mermaid_frontend import SAMPLES, SOURCE_ONLY_SAMPLES

from diagram_pptx import compile_diagram, parse_mermaid
from diagram_pptx.importers import import_mermaid_svg
from diagram_pptx.official import mmdc_version
from diagram_pptx.render import PythonPptxRenderer
from diagram_pptx.scene import Box, Point, SceneConnector, SceneShape, SceneText


def test_mermaid_svg_importer_keeps_native_geometry_and_ignores_images() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">
  <style>#test .messageLine0 { stroke: #333333; }</style>
  <script>alert(1)</script>
  <image href="https://example.com/a.png" width="10" height="10"/>
  <g id="flowchart-A-0" class="node">
    <rect x="10" y="10" width="60" height="30" fill="#abcdef"/>
    <text x="40" y="30">Alpha</text>
  </g>
  <path id="L-A-B-0" class="flowchart-link" d="M 70 25 L 130 25"
        fill="none" stroke="#123456"/>
  <g id="flowchart-B-1" class="node">
    <rect x="130" y="10" width="60" height="30"/>
  </g>
  <polygon id="custom-C" points="50,60 90,90 10,90" fill="#abcdef"/>
  <g id="flowchart-D-2" class="node" transform="translate(100,50)">
    <path class="label-container" d="M -10 -10 h 20 v 20 h -20 z"/>
  </g>
  <line data-id="i0" data-from="A" data-to="B"
        class="messageLine0" x1="20" y1="50" x2="180" y2="50"
        stroke="none" style="fill: none;"/>
</svg>
""",
        kind="flowchart",
    )

    assert any(isinstance(item, SceneShape) for item in scene.elements)
    assert any(isinstance(item, SceneConnector) for item in scene.elements)
    assert any(
        isinstance(item, SceneShape) and item.shape == "custom" and item.points
        for item in scene.elements
    )
    relative_path = next(
        item for item in scene.elements if isinstance(item, SceneShape) and item.semantic_id == "D"
    )
    assert relative_path.box.width == pytest.approx(20)
    assert relative_path.box.height == pytest.approx(20)
    assert all("image" not in item.role for item in scene.elements)
    message = next(
        item
        for item in scene.elements
        if isinstance(item, SceneConnector) and item.semantic_id == "event-0"
    )
    assert message.style.line == "#333333"
    assert (message.source_id, message.target_id) == ("A", "B")
    alpha_shapes = [
        item for item in scene.elements if isinstance(item, SceneShape) and item.semantic_id == "A"
    ]
    assert alpha_shapes[-1].text == "Alpha"
    assert not any(
        isinstance(item, SceneText) and item.semantic_id == "A" for item in scene.elements
    )

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = PythonPptxRenderer().render(
        scene,
        slide=slide,
        bounds=(1, 1, 8, 4),
    )
    assert result.group_shape is not None


def test_mermaid_svg_ignores_empty_label_rectangles() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 50">
  <g id="flowchart-A-0" class="node">
    <rect x="10" y="10" width="80" height="30"/>
    <g class="label"><rect/><text x="50" y="30">Alpha</text></g>
  </g>
</svg>
""",
        kind="flowchart",
    )

    shapes = [item for item in scene.elements if isinstance(item, SceneShape)]
    assert len(shapes) == 1
    assert shapes[0].text == "Alpha"


def test_mermaid_svg_puts_kanban_card_text_in_the_card_shape() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 220 80">
  <g id="flowchart-task1-0" class="node items">
    <rect class="label-container" x="10" y="10" width="200" height="50" rx="4"/>
    <text class="label" x="110" y="40" text-anchor="middle">Research</text>
  </g>
</svg>
""",
        kind="kanban",
    )

    card = next(item for item in scene.elements if isinstance(item, SceneShape))
    assert card.text == "Research"
    assert not any(isinstance(item, SceneText) for item in scene.elements)


def test_mermaid_svg_deduplicates_journey_labels_into_their_shapes() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 240 80">
  <rect class="journey-section" x="10" y="10" width="220" height="50" rx="4"/>
  <text x="120" y="40" text-anchor="middle">Discover</text>
  <text class="journey-section" x="120" y="40"
        text-anchor="middle" font-size="18">Discover</text>
</svg>
""",
        kind="journey",
    )

    section = next(item for item in scene.elements if isinstance(item, SceneShape))
    assert section.text == "Discover"
    assert section.style.font_size == 18
    assert not any(isinstance(item, SceneText) for item in scene.elements)


def test_mermaid_svg_puts_timeline_labels_in_their_shapes() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 240 100">
  <rect class="timeline-node eventWrapper" x="10" y="20"
        width="220" height="50" rx="4"/>
  <text class="timeline-node eventWrapper" x="120" y="68"
        text-anchor="middle" font-size="16">Public beta</text>
</svg>
""",
        kind="timeline",
    )

    event = next(item for item in scene.elements if isinstance(item, SceneShape))
    assert event.text == "Public beta"
    assert not any(isinstance(item, SceneText) for item in scene.elements)


def test_mermaid_svg_samples_elliptical_arc_paths() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 40 40">
  <path d="M 20 30 m -10 0
           a 10 10 0 1 0 20 0
           a 10 10 0 1 0 -20 0"
        fill="#5353FF" fill-opacity="0.1"
        stroke="#5353FF" stroke-opacity="0.95"/>
</svg>
""",
        kind="venn",
    )

    circle = next(item for item in scene.elements if isinstance(item, SceneShape))
    assert circle.box.width == pytest.approx(20, abs=0.2)
    assert circle.box.height == pytest.approx(20, abs=0.2)
    assert len(circle.points) >= 24
    assert circle.style.fill == "#5353FF1A"
    assert circle.style.line == "#5353FFF2"


def test_mermaid_svg_preserves_text_metrics_rotation_and_tspan_position() -> None:
    scene = import_mermaid_svg(
        """\
<svg id="sample" xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 200">
  <style>
    #sample { font-size: 16px; fill: #123456; }
    #sample .label { font-weight: bold; }
  </style>
  <text class="label" text-anchor="end" dominant-baseline="middle"
        transform="translate(80 100) rotate(-90)">
    <tspan x="0" dy="1em">Rotated label</tspan>
  </text>
</svg>
""",
        kind="quadrant",
    )

    label = next(item for item in scene.elements if isinstance(item, SceneText))
    assert label.align == "right"
    assert label.rotation == pytest.approx(-90)
    assert label.style.font_size == 16
    assert label.style.bold is True
    assert label.style.text == "#123456"
    assert label.box.width > 80
    assert label.box.height == pytest.approx(19.2)


def test_mermaid_svg_enlarges_radar_typography_for_slide_readability() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 700 700">
  <circle class="radarGraticule" cx="350" cy="350" r="300"/>
  <line class="radarAxisLine" x1="350" y1="350" x2="350" y2="50"/>
  <text class="radarAxisLabel" x="350" y="30"
        text-anchor="middle" font-size="12">Speed</text>
  <text class="radarLegendText" x="610" y="90"
        font-size="12">Alpha</text>
  <text class="radarTitle" x="350" y="15"
        text-anchor="middle" font-size="16">Product comparison</text>
</svg>
""",
        kind="radar",
    )

    labels = {item.text: item for item in scene.elements if isinstance(item, SceneText)}
    outer = next(item for item in scene.elements if isinstance(item, SceneShape))
    assert outer.box.width == pytest.approx(516)
    assert labels["Speed"].style.font_size == 30
    assert labels["Alpha"].style.font_size == 22
    assert labels["Product comparison"].style.font_size == 26
    assert labels["Speed"].box.width == pytest.approx(84.6)
    assert labels["Product comparison"].box.center.x == pytest.approx(labels["Speed"].box.center.x)


def test_mermaid_svg_compacts_quadrant_plot_and_enlarges_labels() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 500 500">
  <rect class="quadrant" x="50" y="50" width="200" height="200"/>
  <rect class="quadrant" x="250" y="50" width="200" height="200"/>
  <rect class="quadrant" x="50" y="250" width="200" height="200"/>
  <rect class="quadrant" x="250" y="250" width="200" height="200"/>
  <line class="border" x1="50" y1="50" x2="450" y2="50"/>
  <circle class="data-point" cx="150" cy="150" r="5"/>
  <text class="data-point" x="150" y="170"
        text-anchor="middle" font-size="12">Automate</text>
  <text class="labels" x="150" y="480"
        text-anchor="middle" font-size="16">Low effort</text>
  <text class="title" x="250" y="25"
        text-anchor="middle" font-size="20">Portfolio</text>
</svg>
""",
        kind="quadrant",
    )

    quadrants = [
        item
        for item in scene.elements
        if isinstance(item, SceneShape) and "quadrant" in item.classes
    ]
    labels = {item.text: item for item in scene.elements if isinstance(item, SceneText)}
    border = next(item for item in scene.elements if isinstance(item, SceneConnector))
    assert sum(item.box.width for item in quadrants[:2]) == pytest.approx(344)
    assert border.points[0].x == pytest.approx(min(item.box.x for item in quadrants))
    assert labels["Automate"].style.font_size == pytest.approx(22.5)
    assert labels["Low effort"].style.font_size == 20
    assert labels["Portfolio"].style.font_size == 25


def test_mermaid_svg_stacks_and_enlarges_sankey_name_value_labels() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 600 400">
  <rect class="node" x="20" y="20" width="10" height="300"/>
  <text class="node-labels" x="36" y="170"
        font-size="14">Visitors 100</text>
</svg>
""",
        kind="sankey",
    )

    label = next(item for item in scene.elements if isinstance(item, SceneText))
    assert label.text == "Visitors\n100"
    assert label.style.font_size == 22
    assert label.box.height == pytest.approx(52.8)


def test_mermaid_svg_composes_nested_transforms_and_nested_svg_viewbox() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 200">
  <g transform="translate(100 50)">
    <g transform="scale(2)">
      <rect x="5" y="5" width="10" height="10" fill="#abcdef"/>
    </g>
    <svg x="20" y="30" width="20" height="10" viewBox="0 0 80 40">
      <rect width="80" height="40" fill="#123456"/>
    </svg>
  </g>
</svg>
""",
        kind="architecture",
    )

    shapes = [item for item in scene.elements if isinstance(item, SceneShape)]
    scaled, viewport = shapes
    assert scaled.box == Box(0, 0, 20, 20)
    assert viewport.box == Box(10, 20, 20, 10)


def test_mermaid_svg_skips_hidden_and_transparent_helper_geometry() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <clipPath id="clip"><rect width="90" height="90"/></clipPath>
  <rect width="80" height="80" style="display: none"/>
  <path d="M 10 10 A 20 20 0 0 1 50 50"
        style="fill: transparent; fill-opacity: 0"/>
  <circle cx="50" cy="50" r="10" fill="#abcdef"/>
</svg>
""",
        kind="venn",
    )

    shapes = [item for item in scene.elements if isinstance(item, SceneShape)]
    assert len(shapes) == 1
    assert shapes[0].shape == "ellipse"


def test_mermaid_svg_preserves_dom_order_and_multiline_tspans() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">
  <rect width="200" height="100" fill="#ffffff"/>
  <line x1="10" y1="50" x2="190" y2="50" stroke="#333333"/>
  <text x="100" y="42" text-anchor="middle">
    <tspan x="100" dy="0">First</tspan>
    <tspan x="100" dy="16">Second</tspan>
  </text>
</svg>
""",
        kind="zenuml",
    )

    background, connector, label = scene.elements
    assert background.z_index < connector.z_index < label.z_index
    assert isinstance(label, SceneText)
    assert label.text == "First\nSecond"
    assert label.box.height >= 33


def test_mermaid_svg_resolves_compound_class_text_anchor_and_corner_radius() -> None:
    scene = import_mermaid_svg(
        """\
<svg id="sample" xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">
  <style>
    #sample .label { text-anchor: end; }
    #sample .label.cause { text-anchor: middle; font-weight: bold; }
  </style>
  <rect x="20" y="20" width="100" height="40" rx="4" fill="#ffffff"/>
  <text class="label cause" x="70" y="40">Centered</text>
</svg>
""",
        kind="ishikawa",
    )

    shape = next(item for item in scene.elements if isinstance(item, SceneShape))
    label = next(item for item in scene.elements if isinstance(item, SceneText))
    assert shape.metadata["corner_radius_ratio"] == pytest.approx(0.1)
    assert label.align == "center"
    assert label.style.bold is True


def test_mermaid_svg_preserves_rotated_rectangle_geometry() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <rect x="-20" y="-10" width="40" height="20"
        transform="translate(50 50) rotate(-45)" fill="#ffffff"/>
</svg>
""",
        kind="gitgraph",
    )

    shape = next(item for item in scene.elements if isinstance(item, SceneShape))
    assert shape.box.width == pytest.approx(40)
    assert shape.box.height == pytest.approx(20)
    assert shape.rotation == pytest.approx(-45)


def test_mermaid_svg_applies_css_transform_origin() -> None:
    scene = import_mermaid_svg(
        """\
<svg id="sample" xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <style>#sample .milestone { transform: rotate(45deg); }</style>
  <rect x="0" y="0" width="1" height="1" fill="#ffffff"/>
  <rect class="milestone" x="40" y="40" width="20" height="20"
        transform-origin="50px 50px" fill="#abcdef"/>
</svg>
""",
        kind="gantt",
    )

    shape = next(
        item
        for item in scene.elements
        if isinstance(item, SceneShape) and "milestone" in item.classes
    )
    assert shape.box.center == Point(50, 50)
    assert shape.rotation == pytest.approx(45)


def test_mermaid_svg_inherits_group_presentation_attributes() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g fill="none" stroke="#333333">
    <path d="M 10 10 V 90 H 90"/>
  </g>
</svg>
""",
        kind="gantt",
    )

    assert len(scene.elements) == 1
    path = scene.elements[0]
    assert isinstance(path, SceneConnector)
    assert path.style.fill == "none"
    assert path.style.line == "#333333"


def test_mermaid_svg_applies_xhtml_label_css_to_foreign_object_text() -> None:
    scene = import_mermaid_svg(
        """\
<svg id="sample" xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 50">
  <style>#sample .root span { color: #ffffff; font-weight: bold; }</style>
  <g class="root">
    <foreignObject x="10" y="10" width="80" height="30">
      <div xmlns="http://www.w3.org/1999/xhtml">
        <span class="nodeLabel"><p>Product</p></span>
      </div>
    </foreignObject>
  </g>
</svg>
""",
        kind="mindmap",
    )

    label = next(item for item in scene.elements if isinstance(item, SceneText))
    assert label.style.text == "#FFFFFF"
    assert label.style.fill is None
    assert label.style.bold is True


def test_sequence_self_message_is_enlarged_and_labels_are_offset() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 180 120">
  <text class="messageText" x="50" y="32">Prepare</text>
  <path data-id="i0" data-from="A" data-to="A" class="messageLine0"
        d="M 50 60 L 80 60 L 80 78 L 50 78"
        fill="none" stroke="#333333" marker-end="url(#arrow)"/>
  <text class="sequenceNumber" x="50" y="63">1</text>
</svg>
""",
        kind="sequence",
    )

    connector = next(item for item in scene.elements if isinstance(item, SceneConnector))
    label = next(
        item
        for item in scene.elements
        if isinstance(item, SceneText) and "messageText" in item.classes
    )
    number = next(
        item
        for item in scene.elements
        if isinstance(item, SceneText) and "sequenceNumber" in item.classes
    )

    assert connector.metadata["readable_self_message"] is True
    assert max(point.x for point in connector.points) - connector.points[0].x >= 80
    assert connector.points[-1].y - connector.points[0].y >= 30
    assert label.box.center.x > connector.points[0].x
    assert number.box.center.x < connector.points[0].x
    assert number.box.center.y < connector.points[0].y


def test_mermaid_svg_snaps_marker_margin_to_target_boundary() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 60">
  <g id="flowchart-A-0" class="node">
    <rect x="10" y="15" width="50" height="30"/>
  </g>
  <path id="L_A_B_0" class="flowchart-link"
        d="M 60 30 L 126 30" fill="none" marker-end="url(#arrow)"/>
  <g id="flowchart-B-1" class="node">
    <rect x="130" y="15" width="50" height="30"/>
  </g>
</svg>
""",
        kind="flowchart",
    )

    connector = next(item for item in scene.elements if isinstance(item, SceneConnector))
    assert connector.points[-1].x == pytest.approx(120.0)
    assert connector.points[-1].y == pytest.approx(15.0)


def test_class_svg_recovers_solid_frames_source_arrow_and_label_background() -> None:
    scene = import_mermaid_svg(
        """\
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 240 80">
  <g id="classId-Repository-0" class="node">
    <rect x="10" y="10" width="70" height="60"
          fill="#ECECFF" stroke="#9370DB" stroke-dasharray="0 0"/>
  </g>
  <path id="id_Repository_SqlRepository_1"
        class="relation edge-pattern-dashed"
        d="M 90 40 L 150 40" fill="none" stroke="#333333"
        stroke-dasharray="3" marker-start="url(#extension)"/>
  <g class="edgeLabels">
    <g class="edgeLabel">
      <foreignObject x="95" y="28" width="50" height="24">
        <div xmlns="http://www.w3.org/1999/xhtml">implements</div>
      </foreignObject>
    </g>
  </g>
  <g id="classId-SqlRepository-1" class="node">
    <rect x="160" y="10" width="70" height="60"
          fill="#ECECFF" stroke="#9370DB" stroke-dasharray="0 0"/>
  </g>
</svg>
""",
        kind="class",
    )

    frames = [item for item in scene.elements if isinstance(item, SceneShape)]
    assert all(item.style.dash is None for item in frames)
    connector = next(item for item in scene.elements if isinstance(item, SceneConnector))
    assert connector.style.dash == "dash"
    assert connector.start_marker == "triangle"
    assert (connector.source_id, connector.target_id) == (
        "Repository",
        "SqlRepository",
    )
    assert connector.points[0].x == pytest.approx(70.0)
    assert connector.points[-1].x == pytest.approx(150.0)
    label = next(
        item for item in scene.elements if isinstance(item, SceneText) and item.text == "implements"
    )
    assert label.role == "edge.label"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = PythonPptxRenderer().render(
        scene,
        slide=slide,
        bounds=(1, 1, 8, 4),
        group=False,
    )
    relation_parts = result.element_parts[connector.semantic_id]
    assert len(relation_parts) == 1
    assert 'type="triangle"' in relation_parts[0]._element.xml
    assert not any(shape.name.endswith(":marker-start") for shape in relation_parts)


@pytest.mark.official
@pytest.mark.parametrize("kind", list(SAMPLES))
def test_pinned_official_backend_compiles_all_families(
    kind: str,
    tmp_path: Path,
) -> None:
    executable = shutil.which("mmdc")
    if executable is None:
        pytest.skip("mmdc is not installed")
    assert (mmdc_version(executable) or "").startswith("11.16.")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        parse_mermaid(SAMPLES[kind]),
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="official",
        strict=True,
        timeout=60,
    )
    output = tmp_path / f"{kind}.pptx"
    presentation.save(output)

    assert result.backend_used == "official"
    assert result.mermaid_version.startswith("11.16.")
    assert result.scene is not None and result.scene.elements
    assert len(slide.shapes) == 1
    if kind == "sequence":
        assert len(result.nested_group_shapes) == 2
        actor_parts = [item for item in result.scene.elements if "actor-man" in item.classes]
        assert {item.metadata.get("composite_group") for item in actor_parts} == {
            "sequence-actor:actor-top:0",
            "sequence-actor:actor-bottom:0",
        }
        actor_lines = [item for item in actor_parts if isinstance(item, SceneConnector)]
        assert all(item.source_id is None and item.target_id is None for item in actor_lines)
    if kind == "state":
        assert len(result.nested_group_shapes) == 1
        assert result.nested_group_shapes[0].name == "diagram:state:pseudostate:root_end"
        end_parts = [
            item
            for item in result.scene.elements
            if isinstance(item, SceneShape) and item.semantic_id == "root_end"
        ]
        assert len(end_parts) == 4
        assert {item.metadata.get("composite_group") for item in end_parts} == {
            "state-pseudostate:root_end"
        }
    with ZipFile(output) as archive:
        xml = archive.read("ppt/slides/slide1.xml").decode()
    assert "<p:grpSp>" in xml
    if kind == "sequence":
        assert xml.count("<p:grpSp>") >= 3
    if kind == "state":
        assert xml.count("<p:grpSp>") >= 2
    assert "<a:blip" not in xml


@pytest.mark.official
@pytest.mark.parametrize("kind", list(SOURCE_ONLY_SAMPLES))
def test_pinned_official_backend_accepts_every_registered_source_family(
    kind: str,
) -> None:
    executable = shutil.which("mmdc")
    if executable is None:
        pytest.skip("mmdc is not installed")

    document = parse_mermaid(SOURCE_ONLY_SAMPLES[kind], strict=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    result = compile_diagram(
        document,
        slide=slide,
        bounds=(0.7, 0.7, 11.9, 6.1),
        backend="official",
        strict=True,
        timeout=60,
    )

    assert result.backend_used == "official"
    assert result.scene.kind == kind
    assert result.scene.elements
    assert len(slide.shapes) == 1
