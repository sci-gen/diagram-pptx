"""Render a :class:`DrawingScene` as editable native python-pptx shapes."""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from itertools import pairwise
from math import hypot
from typing import Any
from unicodedata import east_asian_width

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from ..scene import (
    Box,
    DrawingScene,
    Point,
    SceneConnector,
    SceneContainer,
    SceneElement,
    SceneShape,
    SceneText,
)
from ..styles import DiagramTheme, ElementStyle, normalize_color


@dataclass(frozen=True, slots=True)
class RenderTheme:
    """Compatibility facade for the pre-alpha fixed renderer theme."""

    node_fill: str = "#EAF2FF"
    node_line: str = "#3167A5"
    node_text: str = "#16324F"
    edge_line: str = "#536273"
    edge_label_fill: str = "#FFFFFF"
    group_line: str = "#8FA4B8"
    group_text: str = "#52687A"
    font_family: str = "Aptos"
    node_font_size: float = 15.0
    edge_font_size: float = 11.0
    group_font_size: float = 11.0
    line_width: float = 1.4

    def to_diagram_theme(self) -> DiagramTheme:
        return DiagramTheme(
            roles={
                "node.default": ElementStyle(
                    fill=self.node_fill,
                    line=self.node_line,
                    text=self.node_text,
                    font_family=self.font_family,
                    font_size=self.node_font_size,
                    line_width=self.line_width,
                ),
                "edge.default": ElementStyle(
                    line=self.edge_line,
                    text=self.edge_line,
                    label_fill=self.edge_label_fill,
                    font_family=self.font_family,
                    font_size=self.edge_font_size,
                    line_width=self.line_width,
                ),
                "group.default": ElementStyle(
                    line=self.group_line,
                    text=self.group_text,
                    font_family=self.font_family,
                    font_size=self.group_font_size,
                ),
            }
        )


@dataclass(slots=True)
class RenderResult:
    """Native objects created for one scene."""

    element_shapes: dict[str, Any] = field(default_factory=dict)
    element_parts: dict[str, list[Any]] = field(default_factory=dict)
    group_shape: Any | None = None
    top_level_shapes: list[Any] = field(default_factory=list)
    connectors: list[Any] = field(default_factory=list)
    edge_label_shapes: list[Any] = field(default_factory=list)
    group_shapes: list[Any] = field(default_factory=list)
    nested_group_shapes: list[Any] = field(default_factory=list)

    def add(self, semantic_id: str, shape: Any, *, primary: bool = False) -> None:
        self.element_parts.setdefault(semantic_id, []).append(shape)
        if primary or semantic_id not in self.element_shapes:
            self.element_shapes[semantic_id] = shape

    @property
    def node_shapes(self) -> dict[str, Any]:
        return self.element_shapes

    @property
    def shapes(self) -> list[Any]:
        """All editable child shapes, excluding the optional outer group."""

        return [shape for parts in self.element_parts.values() for shape in parts]


@dataclass(frozen=True, slots=True)
class _Transform:
    logical_left: float
    logical_top: float
    scale: float
    font_scale: float
    stroke_scale: float
    offset_x: float
    offset_y: float

    def point(self, point: Point) -> tuple[Any, Any]:
        return (
            Inches(self.offset_x + (point.x - self.logical_left) * self.scale),
            Inches(self.offset_y + (point.y - self.logical_top) * self.scale),
        )

    def box(self, box: Box) -> tuple[Any, Any, Any, Any]:
        left, top = self.point(Point(box.x, box.y))
        return (
            left,
            top,
            Inches(box.width * self.scale),
            Inches(box.height * self.scale),
        )


class PythonPptxRenderer:
    """Render a :class:`DrawingScene` as editable ``python-pptx`` objects.

    Most users should call :func:`diagram_pptx.render_mermaid` or
    :func:`diagram_pptx.compile_diagram`. This class is the advanced extension
    point for callers that already produce a positioned drawing scene.
    """

    DEFAULT_SHAPE_REGISTRY = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "ellipse": MSO_SHAPE.OVAL,
        "hexagon": MSO_SHAPE.HEXAGON,
        "stadium": MSO_SHAPE.ROUNDED_RECTANGLE,
        "subprocess": MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
        "cylinder": MSO_SHAPE.CAN,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "custom": MSO_SHAPE.RECTANGLE,
    }

    def __init__(
        self,
        *,
        theme: RenderTheme | None = None,
        shape_registry: dict[str, Any] | None = None,
        inner_padding: float = 0.12,
    ) -> None:
        self.legacy_theme = theme
        self.shape_registry = dict(self.DEFAULT_SHAPE_REGISTRY)
        if shape_registry:
            self.shape_registry.update(shape_registry)
        self.inner_padding = inner_padding

    def render(
        self,
        scene: DrawingScene,
        *,
        target: Any | None = None,
        slide: Any | None = None,
        bounds: tuple[float, float, float, float],
        group: bool = True,
        group_name: str | None = None,
        **_: Any,
    ) -> RenderResult:
        """Render one scene into an existing slide.

        Args:
            scene: Renderer-neutral positioned shapes, connectors, and text.
            target: Preferred name for the destination ``python-pptx`` slide.
            slide: Compatibility alias for ``target``.
            bounds: ``(left, top, width, height)`` in inches.
            group: Move all generated objects into one editable outer group.
            group_name: Optional PowerPoint shape name for the outer group.

        Returns:
            A :class:`RenderResult` containing native shape handles.
        """

        target_slide = target if target is not None else slide
        if target_slide is None:
            raise TypeError("render() requires target=<slide> or slide=<slide>")
        if bounds[2] <= 0 or bounds[3] <= 0:
            raise ValueError("bounds width and height must be positive")
        if not scene.elements:
            return RenderResult()

        transform = self._build_transform(scene, bounds)
        result = RenderResult()
        pending_connections: list[tuple[SceneConnector, list[Any]]] = []
        rendered_by_scene_id: dict[str, list[Any]] = {}
        scene_shapes = {
            item.semantic_id: item
            for item in scene.elements
            if isinstance(item, (SceneShape, SceneContainer))
        }

        for element in sorted(
            enumerate(scene.elements), key=lambda item: (item[1].z_index, item[0])
        ):
            item = element[1]
            item.metadata.setdefault("_diagram_kind", scene.kind)
            rendered_parts: list[Any] = []
            if isinstance(item, SceneContainer):
                parts = self._render_container(target_slide, item, transform)
                rendered_parts.extend(parts)
                result.group_shapes.extend(parts)
                for index, shape in enumerate(parts):
                    result.add(item.semantic_id, shape, primary=index == 0)
            elif isinstance(item, SceneConnector):
                parts, label = self._render_connector(
                    target_slide,
                    item,
                    transform,
                )
                rendered_parts.extend(parts)
                result.connectors.extend(parts)
                for index, shape in enumerate(parts):
                    result.add(item.semantic_id, shape, primary=index == 0)
                if label is not None:
                    rendered_parts.append(label)
                    result.edge_label_shapes.append(label)
                    result.add(item.semantic_id, label)
                pending_connections.append((item, parts))
            elif isinstance(item, SceneShape):
                shape = self._render_shape(target_slide, item, transform)
                rendered_parts.append(shape)
                result.add(item.semantic_id, shape, primary=True)
            elif isinstance(item, SceneText):
                shape = self._render_text(target_slide, item, transform)
                rendered_parts.append(shape)
                result.add(item.semantic_id, shape)
            for shape in rendered_parts:
                self._disable_theme_effects(shape)
            rendered_by_scene_id[item.id] = rendered_parts

        self._bind_connectors(pending_connections, result, scene_shapes)
        nested_groups, root_nested_groups, nested_child_ids = self._group_composites(
            target_slide,
            scene,
            rendered_by_scene_id,
        )
        result.nested_group_shapes.extend(nested_groups)
        child_shapes = [
            shape
            for shape in _unique_shapes(result.shapes)
            if id(shape._element) not in nested_child_ids
        ]
        child_shapes.extend(root_nested_groups)
        if group and child_shapes:
            group_shape = target_slide.shapes.add_group_shape(child_shapes)
            group_shape.name = group_name or f"diagram:{scene.kind}"
            result.group_shape = group_shape
            result.top_level_shapes = [group_shape]
        else:
            result.top_level_shapes = child_shapes
        return result

    @staticmethod
    def _group_composites(
        slide: Any,
        scene: DrawingScene,
        rendered_by_scene_id: dict[str, list[Any]],
    ) -> tuple[list[Any], list[Any], set[int]]:
        grouped_shapes: dict[str, list[Any]] = {}
        group_names: dict[str, str] = {}
        group_parents: dict[str, str] = {}
        for item in scene.elements:
            group_id = item.metadata.get("composite_group")
            if not group_id:
                continue
            grouped_shapes.setdefault(str(group_id), []).extend(
                rendered_by_scene_id.get(item.id, [])
            )
            group_names.setdefault(
                str(group_id),
                str(
                    item.metadata.get(
                        "composite_group_name",
                        f"diagram:{scene.kind}:composite:{group_id}",
                    )
                ),
            )
            parent_id = item.metadata.get("composite_group_parent")
            if parent_id:
                group_parents.setdefault(str(group_id), str(parent_id))

        def depth(group_id: str, seen: set[str] | None = None) -> int:
            active = set() if seen is None else set(seen)
            if group_id in active:
                raise ValueError(f"Composite group cycle detected at {group_id!r}")
            active.add(group_id)
            parent_id = group_parents.get(group_id)
            return 0 if parent_id is None else 1 + depth(parent_id, active)

        nested_groups: list[Any] = []
        nested_by_id: dict[str, Any] = {}
        nested_child_ids: set[int] = set()
        root_nested_groups: list[Any] = []
        children_by_parent: dict[str, list[Any]] = {}
        for group_id in sorted(grouped_shapes, key=lambda item: (-depth(item), item)):
            children = _unique_shapes(
                [
                    *grouped_shapes[group_id],
                    *children_by_parent.get(group_id, []),
                ]
            )
            if len(children) < 2:
                continue
            nested = slide.shapes.add_group_shape(children)
            nested.name = group_names[group_id]
            nested_groups.append(nested)
            nested_by_id[group_id] = nested
            nested_child_ids.update(id(shape._element) for shape in children)
            parent_id = group_parents.get(group_id)
            if parent_id is None:
                root_nested_groups.append(nested)
            else:
                children_by_parent.setdefault(parent_id, []).append(nested)

        # A malformed parent reference must not make an otherwise valid group
        # disappear from the outer diagram group.
        for group_id, nested in nested_by_id.items():
            parent_id = group_parents.get(group_id)
            if parent_id is not None and parent_id not in nested_by_id:
                root_nested_groups.append(nested)
        return nested_groups, _unique_shapes(root_nested_groups), nested_child_ids

    @staticmethod
    def _build_transform(
        scene: DrawingScene, bounds: tuple[float, float, float, float]
    ) -> _Transform:
        scene.recompute_extents()
        logical_width = max(scene.width, 0.001)
        logical_height = max(scene.height, 0.001)
        scale = min(bounds[2] / logical_width, bounds[3] / logical_height)
        rendered_width = logical_width * scale
        rendered_height = logical_height * scale
        svg_units = scene.metadata.get("coordinate_units") == "svg_px"
        metric_scale = scale * 72.0 if svg_units else scale
        return _Transform(
            logical_left=0.0,
            logical_top=0.0,
            scale=scale,
            font_scale=metric_scale,
            stroke_scale=metric_scale,
            offset_x=bounds[0] + (bounds[2] - rendered_width) / 2,
            offset_y=bounds[1] + (bounds[3] - rendered_height) / 2,
        )

    def _render_container(
        self, slide: Any, item: SceneContainer, transform: _Transform
    ) -> list[Any]:
        left, top, width, height = transform.box(item.box)
        is_sequence_fragment = item.role.startswith("sequence.fragment.")
        outline_type = MSO_SHAPE.RECTANGLE if is_sequence_fragment else MSO_SHAPE.ROUNDED_RECTANGLE
        outline = slide.shapes.add_shape(outline_type, left, top, width, height)
        outline.name = self._shape_name(item)
        if not is_sequence_fragment and len(outline.adjustments):
            outline.adjustments[0] = 0.06
        outline.fill.background()
        self._apply_line(
            outline.line,
            item.style,
            default="#8FA4B8",
            scale=transform.stroke_scale,
        )
        if item.style.dash is None and not is_sequence_fragment:
            self._set_dash(outline.line, "dash")
        label_height = min(Inches(0.35), height)
        if is_sequence_fragment:
            label_width = int(
                min(
                    max(Inches(1.2), Inches(0.13 * len(item.label) + 0.45)),
                    max(Inches(0.1), width * 0.45),
                )
            )
            label = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left,
                top,
                label_width,
                label_height,
            )
            self._apply_fill(label.fill, ElementStyle(fill="#FFFFFF"), default="#FFFFFF")
            self._apply_line(
                label.line,
                item.style,
                default="#697586",
                scale=transform.stroke_scale,
            )
        else:
            label = slide.shapes.add_textbox(
                left + Inches(0.08),
                top + Inches(0.02),
                max(Inches(0.1), width - Inches(0.16)),
                label_height,
            )
        label.name = f"{self._shape_name(item)}:label"
        self._format_text_frame(
            label.text_frame,
            item.label,
            item.style,
            default_color="#52687A",
            default_size=11.0,
            align=PP_ALIGN.LEFT,
            font_scale=transform.font_scale,
            min_size=9.0,
            max_size=18.0,
        )
        return [outline, label]

    def _render_shape(self, slide: Any, item: SceneShape, transform: _Transform) -> Any:
        if item.shape == "custom" and len(item.points) >= 3:
            return self._render_freeform(slide, item, transform)
        left, top, width, height = transform.box(item.box)
        shape_type = self.shape_registry.get(item.shape, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.name = self._shape_name(item)
        shape.rotation = item.rotation
        if item.shape == "rounded_rectangle" and len(shape.adjustments):
            shape.adjustments[0] = max(
                0.01,
                min(0.5, float(item.metadata.get("corner_radius_ratio", 0.08))),
            )
        self._apply_fill(shape.fill, item.style, default="#EAF2FF")
        self._apply_line(
            shape.line,
            item.style,
            default="#3167A5",
            scale=transform.stroke_scale,
        )
        if item.text:
            self._format_text_frame(
                shape.text_frame,
                item.text,
                item.style,
                default_color="#16324F",
                default_size=15.0,
                font_scale=transform.font_scale,
                min_size=12.0,
                max_size=40.0,
            )
        return shape

    def _render_freeform(
        self,
        slide: Any,
        item: SceneShape,
        transform: _Transform,
    ) -> Any:
        vertices = [transform.point(point) for point in item.points]
        builder = slide.shapes.build_freeform(
            start_x=vertices[0][0],
            start_y=vertices[0][1],
            scale=1,
        )
        builder.add_line_segments(vertices[1:], close=True)
        shape = builder.convert_to_shape()
        shape.name = self._shape_name(item)
        self._apply_fill(shape.fill, item.style, default="#EAF2FF")
        self._apply_line(
            shape.line,
            item.style,
            default="#3167A5",
            scale=transform.stroke_scale,
        )
        if item.text:
            self._format_text_frame(
                shape.text_frame,
                item.text,
                item.style,
                default_color="#16324F",
                default_size=15.0,
                font_scale=transform.font_scale,
                min_size=12.0,
                max_size=40.0,
            )
        return shape

    def _render_connector(
        self, slide: Any, item: SceneConnector, transform: _Transform
    ) -> tuple[list[Any], Any | None]:
        connectors: list[Any] = []
        points = _simplify_path(item.points)
        for segment_index, (start, end) in enumerate(pairwise(points)):
            start_x, start_y = transform.point(start)
            end_x, end_y = transform.point(end)
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
            )
            connector.name = f"{self._shape_name(item)}:{segment_index}"
            self._apply_line(
                connector.line,
                item.style,
                default="#536273",
                scale=transform.stroke_scale,
            )
            if segment_index == 0 and item.start_marker:
                self._set_line_end(connector.line, "headEnd", item.start_marker)
            if segment_index == len(points) - 2 and item.directed and item.end_marker:
                self._set_line_end(connector.line, "tailEnd", item.end_marker)
            connectors.append(connector)
        label_shape = None
        if item.label and points:
            raw_label_point = item.metadata.get("label_point")
            midpoint, label_segment = self._path_midpoint_segment(points)
            if isinstance(raw_label_point, (list, tuple)) and len(raw_label_point) == 2:
                midpoint = Point(float(raw_label_point[0]), float(raw_label_point[1]))
            center_x, center_y = transform.point(midpoint)
            label_width_inches = self._connector_label_width(
                item.label,
                item.style,
                font_scale=transform.font_scale,
            )
            label_width = Inches(label_width_inches)
            label_height = Inches(0.36)
            label_top = center_y - Inches(0.18)
            label_left = center_x - label_width // 2
            placement = str(item.metadata.get("label_placement", "auto"))
            segment_start, segment_end = label_segment
            segment_horizontal = abs(segment_end.x - segment_start.x) >= abs(
                segment_end.y - segment_start.y
            )
            segment_length_inches = (
                hypot(
                    segment_end.x - segment_start.x,
                    segment_end.y - segment_start.y,
                )
                * transform.scale
            )
            detached = placement == "above" or (
                placement == "auto"
                and segment_length_inches < label_width_inches + 0.28
                and (len(points) == 2 or segment_horizontal)
            )
            if placement == "above" or (detached and segment_horizontal):
                label_top = center_y - label_height - Inches(0.06)
            elif detached:
                label_left = center_x + Inches(0.06)
            label_shape = slide.shapes.add_textbox(
                label_left,
                label_top,
                label_width,
                label_height,
            )
            label_shape.name = f"{self._shape_name(item)}:label"
            if detached and item.style.label_fill is None:
                label_shape.fill.background()
            else:
                self._apply_fill(
                    label_shape.fill,
                    ElementStyle(fill=item.style.label_fill or "#FFFFFF"),
                    default="#FFFFFF",
                )
            label_shape.line.fill.background()
            self._format_text_frame(
                label_shape.text_frame,
                item.label,
                item.style,
                default_color=item.style.line or "#536273",
                default_size=11.0,
                font_scale=transform.font_scale,
                min_size=9.0,
                max_size=18.0,
            )
        return connectors, label_shape

    def _render_text(self, slide: Any, item: SceneText, transform: _Transform) -> Any:
        left, top, width, height = transform.box(item.box)
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.name = self._shape_name(item)
        shape.rotation = item.rotation
        if item.role == "edge.label":
            self._apply_fill(
                shape.fill,
                ElementStyle(fill=item.style.label_fill or item.style.fill or "#FFFFFF"),
                default="#FFFFFF",
            )
        else:
            shape.fill.background()
        shape.line.fill.background()
        align = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[item.align]
        self._format_text_frame(
            shape.text_frame,
            item.text,
            item.style,
            default_color="#16324F",
            default_size=12.0,
            align=align,
            font_scale=transform.font_scale,
            min_size=9.0,
            max_size=40.0,
        )
        shape.text_frame.margin_left = 0
        shape.text_frame.margin_right = 0
        shape.text_frame.margin_top = 0
        shape.text_frame.margin_bottom = 0
        return shape

    @classmethod
    def _bind_connectors(
        cls,
        pending: list[tuple[SceneConnector, list[Any]]],
        result: RenderResult,
        scene_shapes: dict[str, SceneShape | SceneContainer],
    ) -> None:
        for item, parts in pending:
            if not parts:
                continue
            if not item.metadata.get("bind", False):
                continue
            source_shape = result.element_shapes.get(item.source_id or "")
            target_shape = result.element_shapes.get(item.target_id or "")
            source_scene = scene_shapes.get(item.source_id or "")
            target_scene = scene_shapes.get(item.target_id or "")
            source_site = cls._connection_site(source_scene, item.points[0])
            target_site = cls._connection_site(target_scene, item.points[-1])
            if (
                source_shape is not None
                and source_site is not None
                and hasattr(parts[0], "begin_connect")
            ):
                try:
                    parts[0].begin_connect(source_shape, source_site)
                except (ValueError, AttributeError):
                    pass
            if (
                target_shape is not None
                and target_site is not None
                and hasattr(parts[-1], "end_connect")
            ):
                try:
                    parts[-1].end_connect(target_shape, target_site)
                except (ValueError, AttributeError):
                    pass

    @staticmethod
    def _connection_site(
        shape: SceneShape | SceneContainer | None,
        point: Point,
    ) -> int | None:
        if shape is None:
            return None
        mappings = {
            "rectangle": {"top": 0, "left": 1, "bottom": 2, "right": 3},
            "rounded_rectangle": {"top": 0, "left": 1, "bottom": 2, "right": 3},
            "stadium": {"top": 0, "left": 1, "bottom": 2, "right": 3},
            "subprocess": {"top": 0, "left": 1, "bottom": 2, "right": 3},
            "cylinder": {"top": 0, "left": 1, "bottom": 2, "right": 3},
            "diamond": {"top": 0, "left": 1, "bottom": 2, "right": 3},
        }
        mapping = (
            mappings["rounded_rectangle"]
            if isinstance(shape, SceneContainer)
            else mappings.get(shape.shape)
        )
        if mapping is None:
            return None
        distances = {
            "left": abs(point.x - shape.box.x),
            "top": abs(point.y - shape.box.y),
            "right": abs(point.x - (shape.box.x + shape.box.width)),
            "bottom": abs(point.y - (shape.box.y + shape.box.height)),
        }
        return mapping[min(distances, key=distances.get)]

    def _format_text_frame(
        self,
        text_frame: Any,
        text: str,
        style: ElementStyle,
        *,
        default_color: str,
        default_size: float,
        align: Any = PP_ALIGN.CENTER,
        font_scale: float = 1.0,
        min_size: float = 9.0,
        max_size: float = 28.0,
    ) -> None:
        text_frame.clear()
        text_frame.margin_left = Inches(self.inner_padding)
        text_frame.margin_right = Inches(self.inner_padding)
        text_frame.margin_top = Inches(0.04)
        text_frame.margin_bottom = Inches(0.04)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for index, line in enumerate(text.splitlines() or [""]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.alignment = align
            run = paragraph.add_run()
            run.text = line
            run.font.name = style.font_family or "Aptos"
            fitted_size = (style.font_size or default_size) * font_scale
            run.font.size = Pt(max(min_size, min(max_size, fitted_size)))
            run.font.bold = bool(style.bold)
            run.font.italic = bool(style.italic)
            self._apply_font_color(
                run.font.color,
                style.text or default_color,
                opacity=style.opacity,
            )

    @staticmethod
    def _apply_fill(fill: Any, style: ElementStyle, *, default: str) -> None:
        if normalize_color(style.fill or default) in {"none", "transparent"}:
            fill.background()
            return
        fill.solid()
        PythonPptxRenderer._apply_color(
            fill.fore_color,
            style.fill or default,
            opacity=style.opacity,
        )

    @staticmethod
    def _apply_line(
        line: Any,
        style: ElementStyle,
        *,
        default: str,
        scale: float = 1.0,
    ) -> None:
        if normalize_color(style.line or default) in {"none", "transparent"}:
            line.fill.background()
            return
        PythonPptxRenderer._apply_color(
            line.color,
            style.line or default,
            opacity=style.opacity,
        )
        line.width = Pt(max(0.5, min(4.0, (style.line_width or 1.4) * scale)))
        if style.dash:
            PythonPptxRenderer._set_dash(line, style.dash)

    @staticmethod
    def _apply_font_color(
        color_format: Any,
        value: str,
        *,
        opacity: float | None = None,
    ) -> None:
        PythonPptxRenderer._apply_color(color_format, value, opacity=opacity)

    @staticmethod
    def _apply_color(
        color_format: Any,
        value: str,
        *,
        opacity: float | None = None,
    ) -> None:
        normalized = normalize_color(value)
        theme_map = {
            "background1": MSO_THEME_COLOR.BACKGROUND_1,
            "text1": MSO_THEME_COLOR.TEXT_1,
            "background2": MSO_THEME_COLOR.BACKGROUND_2,
            "text2": MSO_THEME_COLOR.TEXT_2,
            "accent1": MSO_THEME_COLOR.ACCENT_1,
            "accent2": MSO_THEME_COLOR.ACCENT_2,
            "accent3": MSO_THEME_COLOR.ACCENT_3,
            "accent4": MSO_THEME_COLOR.ACCENT_4,
            "accent5": MSO_THEME_COLOR.ACCENT_5,
            "accent6": MSO_THEME_COLOR.ACCENT_6,
            "hyperlink": MSO_THEME_COLOR.HYPERLINK,
            "followed_hyperlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
        }
        if normalized in theme_map:
            color_format.theme_color = theme_map[normalized]
            alpha = opacity
        else:
            match = re.fullmatch(
                r"#?([0-9A-Fa-f]{6})([0-9A-Fa-f]{2})?",
                normalized,
            )
            if not match:
                raise ValueError(f"Expected RGB, RGBA, or PowerPoint theme color, got {value!r}")
            color_format.rgb = RGBColor.from_string(match.group(1).upper())
            rgba_alpha = int(match.group(2), 16) / 255 if match.group(2) is not None else 1.0
            alpha = rgba_alpha * (1.0 if opacity is None else opacity)
        if alpha is not None:
            PythonPptxRenderer._set_color_alpha(color_format, alpha)

    @staticmethod
    def _set_color_alpha(color_format: Any, alpha: float) -> None:
        color_element = color_format._color._xClr
        existing = color_element.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
        )
        if existing is not None:
            color_element.remove(existing)
        transparency = OxmlElement("a:alpha")
        transparency.set("val", str(round(max(0.0, min(1.0, alpha)) * 100000)))
        color_element.append(transparency)

    @staticmethod
    def _set_line_end(line: Any, tag: str, marker: str) -> None:
        marker_type = {
            "arrow": "triangle",
            "triangle": "triangle",
            "diamond": "diamond",
            "oval": "oval",
            "circle": "oval",
        }.get(marker)
        if marker.startswith("cardinality:"):
            marker_type = "triangle" if any(char in marker for char in "{}") else "oval"
        if marker_type is None:
            return
        line_element = line._get_or_add_ln()
        existing = line_element.find(
            f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"
        )
        if existing is not None:
            line_element.remove(existing)
        end = OxmlElement(f"a:{tag}")
        end.set("type", marker_type)
        size = "med" if marker in {"arrow", "triangle", "diamond"} else "sm"
        end.set("w", size)
        end.set("len", size)
        line_element.append(end)

    @staticmethod
    def _disable_theme_effects(shape: Any) -> None:
        """Prevent Office theme shadows/effects on explicitly styled objects."""

        style = shape._element.find(qn("p:style"))
        if style is None:
            return
        effect_ref = style.find(qn("a:effectRef"))
        if effect_ref is not None:
            effect_ref.set("idx", "0")

    @staticmethod
    def _connector_label_width(
        text: str,
        style: ElementStyle,
        *,
        font_scale: float,
    ) -> float:
        """Estimate a compact physical label width for Latin and CJK text."""

        columns = sum(
            2 if east_asian_width(character) in {"W", "F", "A"} else 1 for character in text
        )
        font_points = max(9.0, min(18.0, (style.font_size or 11.0) * font_scale))
        return max(0.42, min(3.0, columns * font_points / 72.0 * 0.55 + 0.18))

    @staticmethod
    def _set_dash(line: Any, value: Any) -> None:
        preset = {
            "solid": "solid",
            "dash": "dash",
            "dashed": "dash",
            "dot": "dot",
            "dotted": "dot",
            "dashdot": "dashDot",
        }.get(str(value).replace("-", "").lower(), "dash")
        line_element = line._get_or_add_ln()
        existing = line_element.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash"
        )
        if existing is not None:
            line_element.remove(existing)
        dash = OxmlElement("a:prstDash")
        dash.set("val", preset)
        line_element.append(dash)

    @staticmethod
    def _path_midpoint(points: list[Point]) -> Point:
        return PythonPptxRenderer._path_midpoint_segment(points)[0]

    @staticmethod
    def _path_midpoint_segment(points: list[Point]) -> tuple[Point, tuple[Point, Point]]:
        if len(points) == 1:
            return points[0], (points[0], points[0])
        lengths = [hypot(end.x - start.x, end.y - start.y) for start, end in pairwise(points)]
        total = sum(lengths)
        if total == 0:
            return points[0], (points[0], points[-1])
        target = total / 2
        traversed = 0.0
        for (start, end), length in zip(pairwise(points), lengths, strict=True):
            if traversed + length >= target:
                ratio = (target - traversed) / length if length else 0
                return (
                    Point(
                        start.x + (end.x - start.x) * ratio,
                        start.y + (end.y - start.y) * ratio,
                    ),
                    (start, end),
                )
            traversed += length
        return points[-1], (points[-2], points[-1])

    @staticmethod
    def _shape_name(item: SceneElement) -> str:
        safe_id = re.sub(r"[^A-Za-z0-9_.:-]+", "-", item.semantic_id)
        kind = re.sub(
            r"[^A-Za-z0-9_.:-]+",
            "-",
            str(item.metadata.get("_diagram_kind", "diagram")),
        )
        return f"diagram:{kind}:{item.role}:{safe_id}"


def _unique_shapes(shapes: list[Any]) -> list[Any]:
    result: list[Any] = []
    seen: set[int] = set()
    for shape in shapes:
        identity = id(shape._element)
        if identity not in seen:
            seen.add(identity)
            result.append(shape)
    return result


def _simplify_path(points: list[Point], *, tolerance: float = 1e-7) -> list[Point]:
    """Remove duplicate and collinear vertices without changing the route."""

    deduplicated: list[Point] = []
    for point in points:
        if (
            not deduplicated
            or abs(point.x - deduplicated[-1].x) > tolerance
            or abs(point.y - deduplicated[-1].y) > tolerance
        ):
            deduplicated.append(point)
    if len(deduplicated) < 3:
        return deduplicated
    simplified = [deduplicated[0]]
    for current, following in zip(deduplicated[1:-1], deduplicated[2:], strict=True):
        previous = simplified[-1]
        cross = (current.x - previous.x) * (following.y - current.y) - (current.y - previous.y) * (
            following.x - current.x
        )
        if abs(cross) > tolerance:
            simplified.append(current)
    simplified.append(deduplicated[-1])
    return simplified
