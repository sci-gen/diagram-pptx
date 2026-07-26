"""Command-line interface for compiling and inspecting diagrams."""

from __future__ import annotations

import argparse
import json
import platform
import shutil
import sys
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .compiler import ColorMapName, ColorMapOptions, compile_diagram
from .diagnostics import DiagramPptxError
from .mermaid import parse_mermaid
from .mermaid_registry import MERMAID_COMPATIBILITY_VERSION, mermaid_support_rows
from .official import find_mmdc, mmdc_version
from .styles import DiagramTheme


def _package_version() -> str:
    try:
        return version("diagram-pptx")
    except PackageNotFoundError:
        return "0.1.0b2"


def _optional_package_version(name: str) -> str | None:
    try:
        return version(name)
    except PackageNotFoundError:
        return None


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="diagram-pptx",
        description=(
            "Compile Mermaid-compatible semantic diagrams into editable native PowerPoint shapes."
        ),
    )
    parser.add_argument("--version", action="version", version=f"%(prog)s {_package_version()}")
    subparsers = parser.add_subparsers(dest="command", required=True)

    render = subparsers.add_parser("render", help="Compile a Mermaid file to a new PPTX")
    render.add_argument("input", type=Path, help="Input .mmd file")
    render.add_argument("output", type=Path, help="Output .pptx file")
    render.add_argument("--title", help="Optional slide title outside the diagram group")
    render.add_argument(
        "--bounds",
        help="left,top,width,height in inches",
    )
    render.add_argument(
        "--position",
        choices=["full", "left", "right", "top", "bottom"],
        help="relative placement preset (default: full)",
    )
    render.add_argument(
        "--relative-bounds",
        help="normalized x,y,width,height with a top-left origin and values in 0..1",
    )
    render.add_argument(
        "--backend",
        choices=["auto", "native", "official"],
        default="native",
        help="layout/geometry backend (default: %(default)s)",
    )
    render.add_argument(
        "--style",
        choices=["native", "official"],
        default="native",
        help="visual preset independent of the backend (default: %(default)s)",
    )
    render.add_argument(
        "--colormap",
        choices=["jet", "viridis", "plasma", "magma"],
        help="continuous color map used by numeric color positions",
    )
    render.add_argument(
        "--primary",
        type=float,
        help="shorthand for ordinary node fill position",
    )
    render.add_argument(
        "--secondary",
        type=float,
        help="shorthand for decision fill, outline, and connector positions",
    )
    render.add_argument("--fill", type=float, help="node fill position in the color map")
    render.add_argument("--line", type=float, help="node outline position in the color map")
    render.add_argument("--text", type=float, help="text position in the color map")
    render.add_argument("--edge", type=float, help="connector position in the color map")
    render.add_argument(
        "--decision-fill",
        type=float,
        help="decision-node fill position in the color map",
    )
    render.add_argument("--theme", type=Path, help="Versioned JSON theme/color map")
    render.add_argument(
        "--source-style",
        choices=["merge", "preserve", "replace"],
        default="merge",
    )
    render.add_argument(
        "--label-background",
        help="connector/message-label fill (CSS color, theme slot, or transparent)",
    )
    render.add_argument(
        "--group",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="wrap all diagram shapes in one editable PowerPoint group (default: true)",
    )
    render.add_argument("--mmdc", help="Path to the Mermaid CLI executable")
    render.add_argument("--strict", action="store_true", help="reject unsupported syntax/version")
    render.add_argument("--timeout", type=float, default=30.0, help="mmdc timeout in seconds")

    inspect = subparsers.add_parser("inspect", help="Inspect the typed Mermaid model")
    inspect.add_argument("input", type=Path, help="Input .mmd file")
    inspect.add_argument("--json", action="store_true", help="emit machine-readable JSON")
    inspect.add_argument("--strict", action="store_true", help="reject unsupported syntax")

    doctor = subparsers.add_parser("doctor", help="Check optional rendering runtimes")
    doctor.add_argument("--json", action="store_true", help="emit machine-readable JSON")
    doctor.add_argument("--mmdc", help="Path to the Mermaid CLI executable")

    support = subparsers.add_parser("support", help="Show Mermaid syntax compatibility")
    support.add_argument("--json", action="store_true", help="emit machine-readable JSON")
    return parser


def _parse_bounds(value: str) -> tuple[float, float, float, float]:
    try:
        parts = tuple(float(part.strip()) for part in value.split(","))
    except ValueError as exc:
        raise argparse.ArgumentTypeError("bounds must contain four numbers") from exc
    if len(parts) != 4 or parts[2] <= 0 or parts[3] <= 0:
        raise argparse.ArgumentTypeError("bounds must be left,top,positive-width,positive-height")
    return parts  # type: ignore[return-value]


def render_file(
    input_path: Path,
    output_path: Path,
    *,
    title: str | None = None,
    bounds: tuple[float, float, float, float] | None = None,
    position: str | None = None,
    relative_bounds: tuple[float, float, float, float] | None = None,
    backend: str = "native",
    style: str = "native",
    theme: DiagramTheme | None = None,
    colors: ColorMapOptions | ColorMapName | None = None,
    label_background: str | None = None,
    source_style: str = "merge",
    group: bool = True,
    mmdc_path: str | None = None,
    strict: bool = False,
    timeout: float = 30.0,
) -> Any:
    """Render one Mermaid file to a new one-slide presentation.

    This CLI adapter creates the presentation and slide; library callers should
    normally call :func:`diagram_pptx.render_mermaid` with an existing slide.
    """

    source = input_path.read_text(encoding="utf-8")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    effective_bounds = bounds
    if title:
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.28), Inches(11.93), Inches(0.55))
        paragraph = title_box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = title
        run.font.name = "Aptos Display"
        run.font.size = Pt(28)
        run.font.bold = True
        if bounds is not None:
            effective_bounds = (
                bounds[0],
                max(bounds[1], 1.05),
                bounds[2],
                min(bounds[3], 5.9),
            )
        elif position is None and relative_bounds is None:
            effective_bounds = (0.7, 1.05, 11.93, 5.9)

    result = compile_diagram(
        parse_mermaid(source, strict=strict),
        slide=slide,
        bounds=effective_bounds,
        position=position,
        relative_bounds=relative_bounds,
        backend=backend,
        style=style,
        theme=theme,
        colors=colors,
        label_background=label_background,
        source_style=source_style,  # type: ignore[arg-type]
        group=group,
        mmdc_path=mmdc_path,
        strict=strict,
        timeout=timeout,
    )
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return result


def _inspect_file(path: Path, *, strict: bool, as_json: bool) -> None:
    document = parse_mermaid(path.read_text(encoding="utf-8"), strict=strict)
    payload = {
        "kind": document.model.kind,
        "is_fully_modeled": document.is_fully_modeled,
        "modeling_rate": document.modeling_rate,
        "required_backend": document.required_backend,
        "raw_statements": document.raw_statements,
        "diagnostics": [item.to_dict() for item in document.diagnostics],
        "model": document.model.to_dict(),
    }
    if as_json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return
    print(f"kind: {payload['kind']}")
    print(f"fully modeled: {payload['is_fully_modeled']}")
    print(f"modeling rate: {payload['modeling_rate']:.1%}")
    print(f"required backend: {payload['required_backend']}")
    print(f"diagnostics: {len(document.diagnostics)}")
    for diagnostic in document.diagnostics:
        location = f" line {diagnostic.line}" if diagnostic.line else ""
        print(f"- {diagnostic.severity} {diagnostic.code}{location}: {diagnostic.message}")


def _doctor(*, mmdc_path: str | None, as_json: bool) -> None:
    executable = find_mmdc(mmdc_path)
    resvg_version = _optional_package_version("resvg_py")
    pillow_version = _optional_package_version("Pillow")
    payload = {
        "diagram_pptx": _package_version(),
        "python": platform.python_version(),
        "python_pptx": version("python-pptx"),
        "mmdc": {
            "available": executable is not None,
            "path": executable,
            "version": mmdc_version(executable) if executable else None,
        },
        "image_export": {
            "svg": True,
            "png": resvg_version is not None,
            "jpeg": resvg_version is not None and pillow_version is not None,
            "resvg_py": resvg_version,
            "pillow": pillow_version,
        },
        "libreoffice": shutil.which("libreoffice") or shutil.which("soffice"),
        "pdftoppm": shutil.which("pdftoppm"),
    }
    if as_json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return
    print(f"diagram-pptx: {payload['diagram_pptx']}")
    print(f"Python: {payload['python']}")
    print(f"python-pptx: {payload['python_pptx']}")
    print(
        "mmdc: "
        + (
            f"{payload['mmdc']['version'] or 'unknown'} ({payload['mmdc']['path']})"
            if payload["mmdc"]["available"]
            else "not found (native backend remains available)"
        )
    )
    print("SVG export: available")
    print(
        "PNG/JPEG export: "
        + (
            f"available (resvg_py {resvg_version}, Pillow {pillow_version})"
            if payload["image_export"]["jpeg"]
            else 'not installed (run `pip install "diagram-pptx[image]"`)'
        )
    )
    print(f"LibreOffice: {payload['libreoffice'] or 'not found'}")
    print(f"pdftoppm: {payload['pdftoppm'] or 'not found'}")


def _support(*, as_json: bool) -> None:
    rows = mermaid_support_rows()
    payload = {
        "mermaid_version": MERMAID_COMPATIBILITY_VERSION,
        "families": rows,
        "summary": {
            "registered": len(rows),
            "official": sum(bool(row["official_backend"]) for row in rows),
            "typed_model": sum(bool(row["typed_model"]) for row in rows),
            "native_backend": sum(bool(row["native_backend"]) for row in rows),
        },
    }
    if as_json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return
    print(f"Mermaid compatibility target: {MERMAID_COMPATIBILITY_VERSION}")
    print("family          typed  native  official")
    for row in rows:
        print(
            f"{str(row['kind']):<15} "
            f"{'yes' if row['typed_model'] else 'source':<6} "
            f"{'yes' if row['native_backend'] else 'no':<7} yes"
        )


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        if args.command == "render":
            try:
                bounds = _parse_bounds(args.bounds) if args.bounds else None
                relative_bounds = (
                    _parse_bounds(args.relative_bounds) if args.relative_bounds else None
                )
            except argparse.ArgumentTypeError as exc:
                parser.error(str(exc))
            theme = DiagramTheme.from_json(args.theme) if args.theme else None
            color_channels = {
                name: value
                for name, value in {
                    "primary": args.primary,
                    "secondary": args.secondary,
                    "fill": args.fill,
                    "line": args.line,
                    "text": args.text,
                    "edge": args.edge,
                    "decision_fill": args.decision_fill,
                }.items()
                if value is not None
            }
            if color_channels and args.colormap is None:
                raise ValueError("Color positions require --colormap")
            colors = (
                {"name": args.colormap, **color_channels} if args.colormap is not None else None
            )
            render_file(
                args.input,
                args.output,
                title=args.title,
                bounds=bounds,
                position=args.position,
                relative_bounds=relative_bounds,
                backend=args.backend,
                style=args.style,
                theme=theme,
                colors=colors,  # type: ignore[arg-type]
                label_background=args.label_background,
                source_style=args.source_style,
                group=args.group,
                mmdc_path=args.mmdc,
                strict=args.strict,
                timeout=args.timeout,
            )
            return 0
        if args.command == "inspect":
            _inspect_file(args.input, strict=args.strict, as_json=args.json)
            return 0
        if args.command == "doctor":
            _doctor(mmdc_path=args.mmdc, as_json=args.json)
            return 0
        if args.command == "support":
            _support(as_json=args.json)
            return 0
    except (DiagramPptxError, OSError, ValueError) as exc:
        print(f"diagram-pptx: error: {exc}", file=sys.stderr)
        return 2
    parser.error(f"Unknown command: {args.command}")
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
