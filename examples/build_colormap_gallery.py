"""Build a three-slide gallery for style presets and continuous color maps."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from diagram_pptx import DiagramTheme, ElementStyle, render_mermaid

SOURCE = """\
flowchart LR
    START[評価を開始] --> CHOICE{優先する観点}
    CHOICE -->|速度| FAST[処理を高速化]
    CHOICE -->|品質| QUALITY[品質を向上]
    CHOICE -->|コスト| COST[コストを最適化]
    FAST --> RESULT[実行方針を決定]
    QUALITY --> RESULT
    COST --> RESULT
"""

VARIANTS = (
    (
        "Jet",
        "jet",
        {
            "fill": 0.56,
            "line": 0.90,
            "edge": 0.18,
            "text": 0.20,
            "decision_fill": 0.70,
        },
    ),
    (
        "Viridis",
        "viridis",
        {
            "fill": 0.82,
            "line": 0.25,
            "edge": 0.48,
            "text": 0.02,
            "decision_fill": 0.98,
        },
    ),
    (
        "Plasma",
        "plasma",
        {
            "fill": 0.82,
            "line": 0.18,
            "edge": 0.40,
            "text": 0.02,
            "decision_fill": 0.98,
        },
    ),
    (
        "Viridis two-color",
        "viridis",
        {
            "primary": 0.82,
            "secondary": 0.18,
        },
    ),
    (
        "Magma inverted",
        "magma",
        {
            "primary": 0.08,
            "secondary": 0.68,
        },
    ),
)


def _add_title(slide, title: str, parameters: dict[str, float]) -> None:
    title_box = slide.shapes.add_textbox(
        Inches(0.72),
        Inches(0.28),
        Inches(11.9),
        Inches(0.52),
    )
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = f"{title} · Native / Official style"
    run.font.name = "Noto Sans CJK JP"
    run.font.size = Pt(35)
    run.font.bold = True

    subtitle = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(0.86),
        Inches(11.8),
        Inches(0.34),
    )
    subtitle_paragraph = subtitle.text_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.LEFT
    subtitle_run = subtitle_paragraph.add_run()
    subtitle_run.text = " · ".join(f"{name}={value:.2f}" for name, value in parameters.items())
    subtitle_run.font.name = "Noto Sans CJK JP"
    subtitle_run.font.size = Pt(16)


def build(output: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    theme = DiagramTheme(
        defaults=ElementStyle(font_family="Noto Sans CJK JP"),
    )
    for title, colormap, parameters in VARIANTS:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_title(slide, title, parameters)
        result = render_mermaid(
            SOURCE,
            slide=slide,
            bounds=(0.72, 1.42, 11.9, 5.35),
            style="official",
            theme=theme,
            colors={"name": colormap, **parameters},
            group=True,
            strict=True,
        )
        if result.backend_used != "native":
            raise RuntimeError("The gallery must demonstrate the Native default backend")
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


if __name__ == "__main__":
    destination = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else Path("artifacts/colormap-examples/three-choice-colormap-gallery.pptx")
    )
    build(destination)
