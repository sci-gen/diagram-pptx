"""Build a Native gallery for actors, namespaces, and composite states."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from diagram_pptx import render_mermaid

FEATURES = {
    "sequence": """\
sequenceDiagram
    autonumber
    actor User
    participant API
    participant Worker
    User->>+API: Submit request
    activate API
    API->>+Worker: Run task
    Worker-->>-API: Result
    deactivate Worker
    API-->>User: Completed
    deactivate API
""",
    "class": """\
classDiagram
    direction LR
    namespace Domain {
        class Order {
            +id
            +submit()
        }
        class Customer {
            +email
        }
    }
    namespace Infrastructure {
        class OrderRepository {
            +save(order)
        }
    }
    Customer --> Order : places
    Order --> OrderRepository : stores
""",
    "state": """\
stateDiagram-v2
    direction LR
    Idle --> Workflow : start
    state Workflow {
        [*] --> Validate
        Validate --> Execute : valid
        Execute --> [*]
    }
    Workflow --> Idle : cancel
    Workflow --> Done : complete
""",
}


def main() -> None:
    output = Path(sys.argv[1] if len(sys.argv) > 1 else "artifacts/native-feature-gallery.pptx")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    for source in FEATURES.values():
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        render_mermaid(
            source,
            slide=slide,
            bounds=(0.7, 0.7, 11.9, 6.1),
            backend="native",
            strict=True,
        )
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


if __name__ == "__main__":
    main()
