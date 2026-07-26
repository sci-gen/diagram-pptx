"""Build the 31-family Mermaid 11.16 editable-shape example gallery.

Run inside the project Docker image:

    python examples/build_mermaid_syntax_gallery.py

The script writes one Mermaid source per family, a 31-slide PPTX, and
LibreOffice-rendered PNG previews used by ``examples/README.md``.
"""

from __future__ import annotations

import json
import shutil
import subprocess
import tempfile
from collections.abc import Mapping
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from diagram_pptx import (
    MERMAID_NATIVE_KINDS,
    MERMAID_SYNTAX_FAMILIES,
    compile_diagram,
    parse_mermaid,
)

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "examples" / "mermaid-syntax"
SOURCE_DIR = OUTPUT / "sources"
IMAGE_DIR = OUTPUT / "images"
DECK_PATH = OUTPUT / "mermaid-11.16-gallery.pptx"
MANIFEST_PATH = OUTPUT / "manifest.json"


EXAMPLES: Mapping[str, tuple[str, str]] = {
    "flowchart": (
        "Flowchart",
        """\
flowchart LR
    request([Request]) --> validate{Valid?}
    validate -->|yes| process[Process]
    validate -->|no| revise[Revise]
    revise --> request
    process --> store[(Database)]
""",
    ),
    "swimlanes": (
        "Swimlanes",
        """\
swimlane-beta LR
  subgraph Customer
    request[Submit request]
    receive[Receive result]
  end
  subgraph Operations
    review{Complete?}
    fulfill[Fulfill]
  end
  request --> review
  review -->|Yes| fulfill --> receive
  review -->|No| request
""",
    ),
    "sequence": (
        "Sequence Diagram",
        """\
sequenceDiagram
    autonumber
    actor User
    participant API
    participant Worker
    User->>API: Submit request
    activate API
    API->>Worker: Run task
    Worker-->>API: Result
    API-->>User: Completed
    deactivate API
""",
    ),
    "class": (
        "Class Diagram",
        """\
classDiagram
    direction LR
    class Order {
      +id
      +submit()
    }
    class Customer {
      +email
    }
    class OrderRepository {
      <<interface>>
      +save(order)
    }
    Customer "1" --> "*" Order : places
    OrderRepository <|.. Order : persists
""",
    ),
    "state": (
        "State Diagram",
        """\
stateDiagram-v2
    direction LR
    [*] --> Draft
    Draft --> Review : submit
    Review --> Draft : revise
    Review --> Approved : approve
    Approved --> [*]
""",
    ),
    "er": (
        "Entity Relationship Diagram",
        """\
erDiagram
    CUSTOMER {
      string id PK
      string email UK
    }
    ORDER {
      int id PK
      string customer_id FK
    }
    PRODUCT {
      int id PK
      string name
    }
    CUSTOMER ||--o{ ORDER : places
    ORDER }o--o{ PRODUCT : contains
""",
    ),
    "journey": (
        "User Journey",
        """\
journey
    title Product onboarding
    section Discover
      Visit site: 4: User
      Compare plans: 3: User
    section Start
      Create account: 5: User
      Invite team: 4: User, Admin
""",
    ),
    "gantt": (
        "Gantt",
        """\
gantt
    title Release plan
    dateFormat YYYY-MM-DD
    axisFormat %m/%d
    todayMarker off
    section Build
      API implementation :a1, 2026-08-01, 5d
      UI integration     :after a1, 4d
    section Verify
      Acceptance tests   :3d
      Release            :milestone, 2026-08-13, 0d
""",
    ),
    "pie": (
        "Pie / Donut",
        """\
pie showData
    title Workload
    "Build" : 42
    "Review" : 24
    "Test" : 20
    "Operate" : 14
""",
    ),
    "quadrant": (
        "Quadrant Chart",
        """\
quadrantChart
    title Initiative portfolio
    x-axis Low effort --> High effort
    y-axis Low impact --> High impact
    Automate: [0.22, 0.82]
    Redesign: [0.78, 0.76]
    Tune: [0.30, 0.35]
    Retire: [0.70, 0.22]
""",
    ),
    "requirement": (
        "Requirement Diagram",
        """\
requirementDiagram
    requirement availability {
      id: "R-01"
      text: "Service remains available"
      risk: high
      verifymethod: test
    }
    element api {
      type: service
      docref: "API specification"
    }
    api - satisfies -> availability
""",
    ),
    "gitgraph": (
        "GitGraph",
        """\
gitGraph LR:
    commit id: "init"
    branch develop
    checkout develop
    commit id: "feature"
    checkout main
    commit id: "hotfix"
    checkout develop
    merge main
    checkout main
    merge develop tag: "v1.0"
""",
    ),
    "c4": (
        "C4 Diagram",
        """\
C4Context
    title Checkout context
    Person(customer, "Customer", "Places an order")
    System(shop, "Online Shop", "Accepts and tracks orders")
    System_Ext(payment, "Payment Provider", "Processes payment")
    Rel(customer, shop, "Uses")
    Rel(shop, payment, "Charges", "HTTPS")
""",
    ),
    "mindmap": (
        "Mindmap",
        """\
mindmap
  root((Product))
    Discover
      Research
      Interview
    Build
      Design
      Implement
    Learn
      Measure
      Improve
""",
    ),
    "timeline": (
        "Timeline",
        """\
timeline
    title Product evolution
    2024 : Prototype
         : First users
    2025 : Public beta
         : Team workspace
    2026 : Automation
         : Enterprise launch
""",
    ),
    "zenuml": (
        "ZenUML",
        """\
zenuml
    title Checkout
    @Actor Customer
    Customer->Store: submit()
    Store->Payment: authorize()
    Payment-->Store: approved
    Store-->Customer: confirmation
""",
    ),
    "sankey": (
        "Sankey",
        """\
sankey-beta
Visitors,Signup,65
Visitors,Exit,35
Signup,Trial,48
Signup,Exit,17
Trial,Paid,26
Trial,Exit,22
""",
    ),
    "xychart": (
        "XY Chart",
        """\
xychart-beta
    title "Weekly throughput"
    x-axis [Mon, Tue, Wed, Thu, Fri]
    y-axis "Tasks" 0 --> 20
    bar [8, 12, 10, 16, 18]
    line [7, 9, 12, 14, 17]
""",
    ),
    "block": (
        "Block Diagram",
        """\
block-beta
    columns 3
    input["Input"]
    process["Transform"]
    output["Output"]
    input --> process
    process --> output
""",
    ),
    "packet": (
        "Packet Diagram",
        """\
packet
    title Request packet
    0-7: "Version"
    8-15: "Flags"
    16-31: "Length"
    32-63: "Request ID"
    64-95: "Payload"
""",
    ),
    "kanban": (
        "Kanban",
        """\
kanban
  backlog[Backlog]
    task1[Research]
    task2[Design]
  doing[In progress]
    task3[Implementation]
  done[Done]
    task4[Prototype]
""",
    ),
    "architecture": (
        "Architecture Diagram",
        """\
architecture-beta
    group cloud(cloud)[Cloud]
    service api(server)[API] in cloud
    service worker(server)[Worker] in cloud
    service db(database)[Database] in cloud
    api:R --> L:worker
    worker:B --> T:db
""",
    ),
    "radar": (
        "Radar",
        """\
radar-beta
    title Product comparison
    axis speed["Speed"], quality["Quality"], cost["Cost"]
    axis support["Support"], scale["Scale"]
    curve a["Alpha"]{85, 75, 55, 70, 90}
    curve b["Beta"]{65, 88, 78, 82, 72}
    max 100
    min 0
""",
    ),
    "eventmodeling": (
        "Event Modeling",
        """\
eventmodeling
    tf 01 ui CartUI
    tf 02 cmd AddItem
    tf 03 evt ItemAdded
    tf 04 rmo CartItems
    tf 05 ui CartUI
""",
    ),
    "treemap": (
        "Treemap",
        """\
treemap-beta
    "Platform"
      "API": 35
      "Data": 25
      "Operations": 20
      "Experience": 20
""",
    ),
    "venn": (
        "Venn",
        """\
venn-beta
    set Product
    set Engineering
    set Operations
    union Product,Engineering
    union Engineering,Operations
""",
    ),
    "ishikawa": (
        "Ishikawa",
        """\
ishikawa-beta
    Delayed release
      Process
        Late review
        Unclear scope
      Technology
        Slow build
        Flaky tests
      People
        Ownership gaps
""",
    ),
    "wardley": (
        "Wardley Map",
        """\
wardley-beta
    title Digital service
    anchor User [0.95, 0.55]
    component Experience [0.80, 0.55]
    component API [0.62, 0.55]
    component Platform [0.42, 0.55]
    component Compute [0.22, 0.55]
    User->Experience
    Experience->API
    API->Platform
    Platform->Compute
""",
    ),
    "cynefin": (
        "Cynefin",
        """\
cynefin-beta
    title Delivery decisions
    clear
      "Standard deployment"
    complicated
      "Capacity planning"
    complex
      "New market"
    chaotic
      "Major outage"
    confusion
      "Unclassified signal"
""",
    ),
    "treeview": (
        "TreeView",
        """\
treeView-beta
    Product
      Application
        Web
        Mobile
      Platform
        API
        Data
      Operations
""",
    ),
    "railroad": (
        "Railroad",
        """\
railroad-ebnf-beta
title "Identifier grammar"
letter = "a" | "b" | "c" ;
digit = "0" | "1" | "2" ;
identifier = letter , { letter | digit } ;
""",
    ),
}


def _add_heading(slide, title: str, kind: str, backend: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.24), Inches(9.8), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(35)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x17, 0x24, 0x3B)

    detail = "Pure Python Native" if backend == "native" else "Official → editable shapes"
    meta_box = slide.shapes.add_textbox(Inches(10.25), Inches(0.31), Inches(2.45), Inches(0.35))
    meta_frame = meta_box.text_frame
    meta_frame.clear()
    meta = meta_frame.paragraphs[0]
    meta.alignment = PP_ALIGN.RIGHT
    meta_run = meta.add_run()
    meta_run.text = f"{kind} · {detail}"
    meta_run.font.name = "Aptos"
    meta_run.font.size = Pt(11)
    meta_run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def _write_sources() -> None:
    SOURCE_DIR.mkdir(parents=True, exist_ok=True)
    for kind, (_, source) in EXAMPLES.items():
        (SOURCE_DIR / f"{kind}.mmd").write_text(source, encoding="utf-8")


def _build_deck() -> list[dict[str, object]]:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    manifest: list[dict[str, object]] = []

    for index, (kind, (title, source)) in enumerate(EXAMPLES.items(), start=1):
        backend = "native" if kind in MERMAID_NATIVE_KINDS else "official"
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_heading(slide, title, kind, backend)
        result = compile_diagram(
            parse_mermaid(source),
            slide=slide,
            bounds=(0.62, 1.02, 12.09, 5.9),
            backend=backend,
            style="official",
            group=True,
            strict=True,
            timeout=60,
        )
        manifest.append(
            {
                "slide": index,
                "kind": kind,
                "title": title,
                "backend": backend,
                "backend_used": result.backend_used,
                "source": f"sources/{kind}.mmd",
                "image": f"images/{kind}.png",
                "editable_shape_count": len(result.shapes),
            }
        )

    OUTPUT.mkdir(parents=True, exist_ok=True)
    presentation.save(DECK_PATH)
    return manifest


def _verify_deck(manifest: list[dict[str, object]]) -> None:
    """Verify that the gallery is complete and contains editable shapes."""
    with ZipFile(DECK_PATH) as package:
        slide_names = sorted(
            name
            for name in package.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        if len(slide_names) != len(manifest):
            raise RuntimeError(f"Expected {len(manifest)} slides, got {len(slide_names)}")
        for slide_name in slide_names:
            slide_xml = package.read(slide_name)
            if b"<a:blip" in slide_xml:
                raise RuntimeError(f"{slide_name} contains a raster image")
            if b"<p:grpSp" not in slide_xml:
                raise RuntimeError(f"{slide_name} has no editable diagram group")


def _render_previews(manifest: list[dict[str, object]]) -> None:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        missing = ", ".join(
            name for name, value in (("LibreOffice", soffice), ("pdftoppm", pdftoppm)) if not value
        )
        raise RuntimeError(f"Preview rendering requires: {missing}")

    IMAGE_DIR.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="diagram-pptx-gallery-") as raw_tmp:
        temporary = Path(raw_tmp)
        profile = temporary / "lo-profile"
        subprocess.run(
            [
                soffice,
                "--headless",
                f"-env:UserInstallation=file://{profile}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(temporary),
                str(DECK_PATH),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=180,
        )
        pdf = temporary / f"{DECK_PATH.stem}.pdf"
        prefix = temporary / "slide"
        subprocess.run(
            [
                pdftoppm,
                "-png",
                "-scale-to-x",
                "1200",
                "-scale-to-y",
                "-1",
                str(pdf),
                str(prefix),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=180,
        )
        pages = sorted(temporary.glob("slide-*.png"))
        if len(pages) != len(manifest):
            raise RuntimeError(f"Expected {len(manifest)} rendered pages, got {len(pages)}")
        for page, item in zip(pages, manifest, strict=True):
            shutil.copyfile(page, IMAGE_DIR / f"{item['kind']}.png")


def main() -> None:
    registered = {family.kind for family in MERMAID_SYNTAX_FAMILIES}
    included = set(EXAMPLES)
    if included != registered:
        raise RuntimeError(
            "Gallery and compatibility registry differ: "
            f"missing={sorted(registered - included)}, extra={sorted(included - registered)}"
        )
    _write_sources()
    manifest = _build_deck()
    _verify_deck(manifest)
    _render_previews(manifest)
    MANIFEST_PATH.write_text(
        json.dumps(
            {
                "mermaid_version": "11.16.0",
                "deck": DECK_PATH.name,
                "examples": manifest,
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    print(f"Wrote {DECK_PATH}")
    print(f"Wrote {len(manifest)} sources and PNG previews under {OUTPUT}")


if __name__ == "__main__":
    main()
