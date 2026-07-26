"""Build one slide per supported Mermaid family for headless visual review."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from diagram_pptx import render_mermaid

EXAMPLES = {
    "flowchart": """\
flowchart LR
A([Request]) --> B{Valid?}
B -->|yes| C[(Database)]
B -->|no| D[Reject]
""",
    "sequence": """\
sequenceDiagram
actor User
participant API
participant DB
User->>API: Save
API->>DB: Insert
DB-->>API: Result
API-->>User: Created
""",
    "class": """\
classDiagram
direction LR
class Repository {
  <<interface>>
  +find(id)
}
class SqlRepository {
  -connection
  +find(id)
}
Repository <|.. SqlRepository : implements
""",
    "er": """\
erDiagram
CUSTOMER {
  string id PK
  string email UK
}
ORDER {
  int id PK
  string customer_id FK
}
CUSTOMER ||--o{ ORDER : places
""",
    "state": """\
stateDiagram-v2
[*] --> Idle
Idle --> Running : start
Running --> Idle : stop
Running --> [*]
""",
}


def main() -> None:
    output = Path(sys.argv[1] if len(sys.argv) > 1 else "diagram-gallery.pptx")
    backend = sys.argv[2] if len(sys.argv) > 2 else "native"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    for source in EXAMPLES.values():
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        render_mermaid(
            source,
            slide=slide,
            bounds=(0.7, 0.7, 11.9, 6.1),
            backend=backend,
            strict=True,
        )
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


if __name__ == "__main__":
    main()
