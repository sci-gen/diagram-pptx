"""Build a gallery showing same-slide diagram placement."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from diagram_pptx import render_mermaid  # noqa: E402

LEFT_FLOW = """\
flowchart TB
    A[左側の入力] --> B{確認}
    B -->|OK| C[左側で完了]
    B -->|再試行| A
"""

RIGHT_FLOW = """\
flowchart TB
    D[右側の入力] --> E{分類}
    E -->|保存| F[(データベース)]
    E -->|通知| G[メッセージ]
"""

TOP_FLOW = """\
flowchart LR
    A[上段の開始] --> B{判定} -->|続行| C[上段の完了]
    B -->|戻る| A
"""

BOTTOM_FLOW = """\
flowchart LR
    D[下段の開始] --> E{検証} -->|合格| F[下段の完了]
    E -->|修正| D
"""

SMALL_FLOWS = (
    """flowchart LR
    A[受付] --> B[分析]
""",
    """flowchart LR
    C[計画] --> D[実行]
""",
    """flowchart LR
    E[確認] --> F[完了]
""",
)


def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    side_by_side = prs.slides.add_slide(blank)
    render_mermaid(
        LEFT_FLOW,
        slide=side_by_side,
        position="left",
        style="official",
        colors={"name": "viridis", "primary": 0.72, "secondary": 0.18},
    )
    render_mermaid(
        RIGHT_FLOW,
        slide=side_by_side,
        position="right",
        style="official",
        colors={"name": "plasma", "primary": 0.72, "secondary": 0.18},
    )

    stacked = prs.slides.add_slide(blank)
    render_mermaid(
        TOP_FLOW,
        slide=stacked,
        position="top",
        style="official",
        colors={"name": "jet", "primary": 0.18, "secondary": 0.82},
    )
    render_mermaid(
        BOTTOM_FLOW,
        slide=stacked,
        position="bottom",
        style="official",
        colors={"name": "magma", "primary": 0.35, "secondary": 0.85},
    )

    normalized = prs.slides.add_slide(blank)
    placements = (
        (0.05, 0.08, 0.38, 0.36),
        (0.57, 0.08, 0.38, 0.36),
        (0.30, 0.56, 0.40, 0.36),
    )
    maps = ("viridis", "plasma", "jet")
    for source, relative_bounds, colormap in zip(
        SMALL_FLOWS,
        placements,
        maps,
        strict=True,
    ):
        render_mermaid(
            source,
            slide=normalized,
            relative_bounds=relative_bounds,
            style="official",
            colors={"name": colormap, "primary": 0.68, "secondary": 0.22},
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


if __name__ == "__main__":
    target = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else Path("artifacts/placement-examples/diagram-placement-gallery.pptx")
    )
    build(target)
    print(target)
